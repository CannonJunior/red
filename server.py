#!/usr/bin/env python3
"""
Simple HTTP server for the Robobrain UI web application.
Serves static files on port 9090.
"""

import os
import sys
import json
import uuid
import time
import urllib.request
import urllib.parse
from http.server import HTTPServer, BaseHTTPRequestHandler
import mimetypes
from pathlib import Path
from datetime import datetime

# Import server decorators
from server_decorators import require_system

# Import debug logger
from debug_logger import debug_log, info_log, error_log, success_log, warning_log

# Import static cache
from static_cache import get_static_cache

# Import Ollama configuration
from ollama_config import ollama_config

# Import server utilities
from server.utils.json_response import send_json_response as send_json_response_util
from server.utils.request_helpers import get_content_type as get_content_type_util, get_request_body as get_request_body_util

# Import route handlers
from server.routes.models import handle_models_api as handle_models_route
from server.routes.chat import handle_chat_api as handle_chat_route
from server.routes.rag import (
    handle_rag_status_api as handle_rag_status_route,
    handle_rag_search_api as handle_rag_search_route,
    handle_rag_query_api as handle_rag_query_route,
    handle_rag_ingest_api as handle_rag_ingest_route,
    handle_rag_documents_api as handle_rag_documents_route,
    handle_rag_analytics_api as handle_rag_analytics_route,
    handle_rag_upload_api as handle_rag_upload_route,
    handle_rag_document_delete_api as handle_rag_document_delete_route
)
from server.routes.cag import (
    handle_cag_status_api as handle_cag_status_route,
    handle_cag_load_api as handle_cag_load_route,
    handle_cag_clear_api as handle_cag_clear_route,
    handle_cag_document_delete_api as handle_cag_document_delete_route,
    handle_cag_query_api as handle_cag_query_route
)
from server.routes.prompts import (
    handle_prompts_list_api as handle_prompts_list_route,
    handle_prompts_create_api as handle_prompts_create_route,
    handle_prompts_detail_api as handle_prompts_detail_route,
    handle_prompts_update_api as handle_prompts_update_route,
    handle_prompts_delete_api as handle_prompts_delete_route,
    handle_prompts_use_api as handle_prompts_use_route,
    handle_prompts_search_api as handle_prompts_search_route
)
from server.routes.search import (
    handle_search_api as handle_search_route,
    handle_search_folders_api as handle_search_folders_route,
    handle_search_create_folder_api as handle_search_create_folder_route,
    handle_search_tags_api as handle_search_tags_route,
    handle_search_add_object_api as handle_search_add_object_route
)
from server.routes.agents import (
    handle_agents_api as handle_agents_route,
    handle_agents_metrics_api as handle_agents_metrics_route,
    handle_agents_detail_api as handle_agents_detail_route
)

# Ollama agents support
try:
    from server.routes.ollama_agents import (
        handle_ollama_agents_api as handle_ollama_agents_route,
        handle_ollama_agent_detail_api as handle_ollama_agent_detail_route,
        handle_ollama_agent_invoke_api as handle_ollama_agent_invoke_route,
        handle_ollama_skills_api as handle_ollama_skills_route,
        handle_ollama_status_api as handle_ollama_status_route
    )
    OLLAMA_AGENTS_AVAILABLE = True
except ImportError as e:
    debug_log(f"⚠️ Ollama agents not available: {e}", "⚠️")
    OLLAMA_AGENTS_AVAILABLE = False

from server.routes.mcp import (
    handle_mcp_servers_api as handle_mcp_servers_route,
    handle_mcp_server_action_api as handle_mcp_server_action_route,
    handle_nlp_parse_task_api as handle_nlp_parse_task_route,
    handle_nlp_capabilities_api as handle_nlp_capabilities_route,
    handle_mcp_metrics_api as handle_mcp_metrics_route
)
from server.routes.visualizations import (
    handle_knowledge_graph_api as handle_knowledge_graph_route,
    handle_performance_dashboard_api as handle_performance_dashboard_route,
    handle_search_results_api as handle_search_results_route
)
from server.routes.opportunities import (
    handle_opportunities_list_api as handle_opportunities_list_route,
    handle_opportunities_create_api as handle_opportunities_create_route,
    handle_opportunities_detail_api as handle_opportunities_detail_route,
    handle_opportunities_update_api as handle_opportunities_update_route,
    handle_opportunities_delete_api as handle_opportunities_delete_route,
    handle_tasks_list_api as handle_tasks_list_route,
    handle_tasks_create_api as handle_tasks_create_route,
    handle_task_get_api as handle_task_get_route,
    handle_task_update_api as handle_task_update_route,
    handle_task_delete_api as handle_task_delete_route,
    handle_task_history_api as handle_task_history_route
)

# Import TODO functionality
try:
    from server.routes.todos import (
        handle_users_list_api as handle_users_list_route,
        handle_users_create_api as handle_users_create_route,
        handle_users_detail_api as handle_users_detail_route,
        handle_users_update_api as handle_users_update_route,
        handle_users_delete_api as handle_users_delete_route,
        handle_lists_list_api as handle_lists_list_route,
        handle_lists_create_api as handle_lists_create_route,
        handle_lists_detail_api as handle_lists_detail_route,
        handle_lists_update_api as handle_lists_update_route,
        handle_lists_delete_api as handle_lists_delete_route,
        handle_lists_share_api as handle_lists_share_route,
        handle_lists_unshare_api as handle_lists_unshare_route,
        handle_lists_shares_api as handle_lists_shares_route,
        handle_shared_lists_api as handle_shared_lists_route,
        handle_todos_list_api as handle_todos_list_route,
        handle_todos_create_api as handle_todos_create_route,
        handle_todos_detail_api as handle_todos_detail_route,
        handle_todos_update_api as handle_todos_update_route,
        handle_todos_delete_api as handle_todos_delete_route,
        handle_todos_complete_api as handle_todos_complete_route,
        handle_todos_archive_api as handle_todos_archive_route,
        handle_todos_today_api as handle_todos_today_route,
        handle_todos_upcoming_api as handle_todos_upcoming_route,
        handle_todos_search_api as handle_todos_search_route,
        handle_todos_parse_api as handle_todos_parse_route,
        handle_tags_list_api as handle_tags_list_route,
        handle_tags_create_api as handle_tags_create_route,
        handle_tags_detail_api as handle_tags_detail_route,
        handle_tags_update_api as handle_tags_update_route,
        handle_tags_delete_api as handle_tags_delete_route,
        handle_todos_history_api as handle_todos_history_route
    )
    TODOS_AVAILABLE = True
    print("✅ TODO system loaded successfully")
except ImportError as e:
    TODOS_AVAILABLE = False
    print(f"⚠️  TODO system not available: {e}")

# Import RAG functionality
try:
    from rag_api import handle_rag_status_request, handle_rag_search_request, handle_rag_query_request, handle_rag_ingest_request, handle_rag_documents_request, handle_rag_analytics_request, handle_rag_document_delete_request, handle_rag_vector_chunks_request
    from knowledge_graph_builder import VectorKnowledgeGraphBuilder
    RAG_AVAILABLE = True
    KNOWLEDGE_GRAPH_AVAILABLE = True
    print("✅ RAG system loaded successfully")
except ImportError as e:
    RAG_AVAILABLE = False
    KNOWLEDGE_GRAPH_AVAILABLE = False
    print(f"⚠️  RAG system not available: {e}")

# Import CAG (Cache-Augmented Generation) functionality
try:
    from cag_api import get_cag_manager
    CAG_AVAILABLE = True
    cag_manager = get_cag_manager()
    print("✅ CAG system loaded successfully")
except ImportError as e:
    CAG_AVAILABLE = False
    cag_manager = None
    print(f"⚠️  CAG system not available: {e}")

# Import Agent system functionality
try:
    # Add current directory to Python path for agent-system imports
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

    from agent_system.mcp.server_manager import ZeroCostMCPServerManager, initialize_default_servers
    from agent_system.agents.agent_config import MojoOptimizedAgentManager
    from agent_system.security.local_security import ZeroCostLocalSecurity
    AGENT_SYSTEM_AVAILABLE = True
    print("✅ Agent system loaded successfully")

    # Initialize agent managers (singleton pattern for 5-user optimization)
    mcp_manager = initialize_default_servers()
    agent_manager = MojoOptimizedAgentManager()
    security_manager = ZeroCostLocalSecurity()
    print("✅ Agent managers initialized")
except ImportError as e:
    AGENT_SYSTEM_AVAILABLE = False
    mcp_manager = None
    agent_manager = None
    security_manager = None
    print(f"⚠️  Agent system not available: {e}")

# Import Search functionality
try:
    from search_api import handle_search_request, handle_folders_request, handle_create_folder_request, handle_tags_request, handle_add_object_request, handle_update_object_request, handle_delete_object_request
    SEARCH_AVAILABLE = True
    print("✅ Search system loaded successfully")
except ImportError as e:
    SEARCH_AVAILABLE = False
    print(f"⚠️  Search system not available: {e}")

# Import Prompts functionality
try:
    from prompts_api import (
        handle_prompts_list_request, handle_prompts_create_request,
        handle_prompts_get_request, handle_prompts_update_request,
        handle_prompts_delete_request, handle_prompts_use_request,
        handle_prompts_search_request
    )
    PROMPTS_AVAILABLE = True
    print("✅ Prompts system loaded successfully")
except ImportError as e:
    PROMPTS_AVAILABLE = False
    print(f"⚠️  Prompts system not available: {e}")


class CustomHTTPRequestHandler(BaseHTTPRequestHandler):
    """Custom handler to serve static files with proper MIME types."""
    
    def do_GET(self):
        """Handle GET requests."""
        try:
            # Handle API routes
            if self.path.startswith('/api/'):
                # RAG API endpoints
                if self.path == '/api/rag/documents' and RAG_AVAILABLE:
                    self.handle_rag_documents_api()
                    return
                elif self.path == '/api/rag/analytics' and RAG_AVAILABLE:
                    self.handle_rag_analytics_api()
                    return
                # Search API endpoints
                elif self.path == '/api/search/folders' and SEARCH_AVAILABLE:
                    self.handle_search_folders_api()
                    return
                elif self.path == '/api/search/tags' and SEARCH_AVAILABLE:
                    self.handle_search_tags_api()
                    return
                # Visualization API endpoints
                elif self.path == '/api/visualizations/knowledge-graph':
                    self.handle_knowledge_graph_api()
                    return
                elif self.path == '/api/visualizations/performance':
                    self.handle_performance_dashboard_api()
                    return
                elif self.path == '/api/visualizations/search-results':
                    self.handle_search_results_api()
                    return
                # Agent System API endpoints
                elif self.path == '/api/agents' and AGENT_SYSTEM_AVAILABLE:
                    self.handle_agents_api()
                    return
                elif self.path == '/api/agents/metrics' and AGENT_SYSTEM_AVAILABLE:
                    self.handle_agents_metrics_api()
                    return
                elif self.path.startswith('/api/agents/') and AGENT_SYSTEM_AVAILABLE:
                    self.handle_agents_detail_api()
                    return
                # Ollama Agent System API endpoints
                elif self.path == '/api/ollama/status':
                    self.handle_ollama_status_api()
                    return
                elif self.path == '/api/ollama/agents':
                    self.handle_ollama_agents_api()
                    return
                elif self.path == '/api/ollama/skills':
                    self.handle_ollama_skills_api()
                    return
                elif self.path.startswith('/api/ollama/agents/'):
                    self.handle_ollama_agent_detail_api()
                    return
                elif self.path == '/api/mcp/servers' and AGENT_SYSTEM_AVAILABLE:
                    self.handle_mcp_servers_api()
                    return
                elif self.path.startswith('/api/mcp/servers/') and AGENT_SYSTEM_AVAILABLE:
                    self.handle_mcp_server_action_api()
                    return
                elif self.path == '/api/mcp/metrics' and AGENT_SYSTEM_AVAILABLE:
                    self.handle_mcp_metrics_api()
                    return
                elif self.path == '/api/nlp/capabilities' and AGENT_SYSTEM_AVAILABLE:
                    self.handle_nlp_capabilities_api()
                    return
                # Prompts API endpoints
                elif self.path == '/api/prompts' and PROMPTS_AVAILABLE:
                    self.handle_prompts_list_api()
                    return
                elif self.path.startswith('/api/prompts/') and PROMPTS_AVAILABLE:
                    self.handle_prompts_detail_api()
                    return
                # TODO API endpoints - specific routes first, then generic
                elif self.path == '/api/todos/users' and TODOS_AVAILABLE:
                    self.handle_users_list_api()
                    return
                elif (self.path == '/api/todos/lists' or self.path.startswith('/api/todos/lists?')) and TODOS_AVAILABLE:
                    self.handle_lists_list_api()
                    return
                elif self.path == '/api/todos/shared' and TODOS_AVAILABLE:
                    self.handle_shared_lists_api()
                    return
                elif self.path == '/api/todos/today' and TODOS_AVAILABLE:
                    self.handle_todos_today_api()
                    return
                elif self.path == '/api/todos/upcoming' and TODOS_AVAILABLE:
                    self.handle_todos_upcoming_api()
                    return
                elif self.path == '/api/todos/search' and TODOS_AVAILABLE:
                    self.handle_todos_search_api()
                    return
                elif self.path == '/api/todos/tags' and TODOS_AVAILABLE:
                    self.handle_tags_list_api()
                    return
                elif (self.path == '/api/todos' or self.path.startswith('/api/todos?')) and TODOS_AVAILABLE:
                    self.handle_todos_list_api()
                    return
                elif self.path.startswith('/api/todos/lists/') and self.path.endswith('/shares') and TODOS_AVAILABLE:
                    self.handle_lists_shares_api()
                    return
                elif self.path.startswith('/api/todos/users/') and TODOS_AVAILABLE:
                    self.handle_users_detail_api()
                    return
                elif self.path.startswith('/api/todos/lists/') and TODOS_AVAILABLE:
                    self.handle_lists_detail_api()
                    return
                elif self.path.startswith('/api/todos/tags/') and TODOS_AVAILABLE:
                    self.handle_tags_detail_api()
                    return
                elif self.path.startswith('/api/todos/') and self.path.endswith('/history') and TODOS_AVAILABLE:
                    self.handle_todos_history_api()
                    return
                elif self.path.startswith('/api/todos/') and TODOS_AVAILABLE:
                    self.handle_todos_detail_api()
                    return
                # Opportunities API endpoints
                elif self.path == '/api/opportunities':
                    self.handle_opportunities_list_api()
                    return
                elif self.path.startswith('/api/opportunities/') and '/tasks' in self.path:
                    # Handle /api/opportunities/{id}/tasks
                    self.handle_tasks_list_api()
                    return
                elif self.path.startswith('/api/tasks/') and self.path.endswith('/history'):
                    # Handle /api/tasks/{id}/history
                    self.handle_task_history_api()
                    return
                elif self.path.startswith('/api/tasks/'):
                    # Handle /api/tasks/{id}
                    self.handle_task_get_api()
                    return
                elif self.path.startswith('/api/opportunities/'):
                    self.handle_opportunities_detail_api()
                    return
                else:
                    self.send_error(404, f"API endpoint not found: {self.path}")
                    return
            
            # Handle root path
            if self.path == '/':
                self.path = '/index.html'
            
            # Remove leading slash and resolve file path
            file_path = self.path.lstrip('/')
            full_path = os.path.join(os.getcwd(), file_path)
            
            debug_log(f"Request: {self.path} -> {file_path}")
            
            # Check if file exists
            if not os.path.exists(full_path) or not os.path.isfile(full_path):
                self.send_error(404, f"File not found: {self.path}")
                return
            
            # Determine content type
            content_type = self.get_content_type(file_path)

            # Try to get from cache first
            static_cache = get_static_cache()
            content, etag = static_cache.get(full_path)

            if content is None:
                # Cache miss - read from disk
                with open(full_path, 'rb') as f:
                    content = f.read()
                # Store in cache
                etag = static_cache.set(full_path, content)
                debug_log(f"Cache MISS: {file_path}", "💾")
            else:
                debug_log(f"Cache HIT: {file_path}", "⚡")

            # Send response
            self.send_response(200)
            self.send_header('Content-Type', content_type)
            self.send_header('Content-Length', str(len(content)))
            self.send_header('ETag', etag)
            self.send_header('Cache-Control', 'public, max-age=3600')
            # Add CORS headers
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
            self.send_header('Access-Control-Allow-Headers', 'Content-Type')
            self.end_headers()
            
            # Send content
            self.wfile.write(content)
            
            debug_log(f"Served {file_path} as {content_type}", "✅")
            
        except Exception as e:
            print(f"❌ Error serving {self.path}: {e}")
            self.send_error(500, f"Internal server error: {e}")
    
    def do_POST(self):
        """Handle POST requests for API endpoints."""
        try:
            # Parse request path
            if self.path == '/api/chat':
                self.handle_chat_api()
            elif self.path == '/api/models':
                self.handle_models_api()
            elif self.path == '/api/rag/status' and RAG_AVAILABLE:
                self.handle_rag_status_api()
            elif self.path == '/api/rag/search' and RAG_AVAILABLE:
                self.handle_rag_search_api()
            elif self.path == '/api/rag/query' and RAG_AVAILABLE:
                self.handle_rag_query_api()
            elif self.path == '/api/rag/ingest' and RAG_AVAILABLE:
                self.handle_rag_ingest_api()
            elif self.path == '/api/rag/upload' and RAG_AVAILABLE:
                self.handle_rag_upload_api()
            elif self.path == '/api/rag/documents' and RAG_AVAILABLE:
                self.handle_rag_documents_api()
            # CAG API endpoints
            elif self.path == '/api/cag/status' and CAG_AVAILABLE:
                self.handle_cag_status_api()
            elif self.path == '/api/cag/load' and CAG_AVAILABLE:
                self.handle_cag_load_api()
            elif self.path == '/api/cag/clear' and CAG_AVAILABLE:
                self.handle_cag_clear_api()
            elif self.path == '/api/cag/query' and CAG_AVAILABLE:
                self.handle_cag_query_api()
            # Agent API endpoints
            elif self.path == '/api/agents' and AGENT_SYSTEM_AVAILABLE:
                self.handle_agents_api()
            elif self.path.startswith('/api/agents/') and AGENT_SYSTEM_AVAILABLE:
                self.handle_agents_detail_api()
            # Ollama Agent API endpoints
            elif self.path == '/api/ollama/agents':
                self.handle_ollama_agents_api()
            elif self.path.startswith('/api/ollama/agents/') and '/invoke' in self.path:
                self.handle_ollama_agent_invoke_api()
            elif self.path.startswith('/api/ollama/agents/'):
                self.handle_ollama_agent_detail_api()
            # Search API endpoints
            elif self.path == '/api/mcp/servers' and AGENT_SYSTEM_AVAILABLE:
                self.handle_mcp_servers_api()
            elif self.path.startswith('/api/mcp/servers/') and AGENT_SYSTEM_AVAILABLE:
                self.handle_mcp_server_action_api()
            elif self.path == '/api/nlp/parse-task' and AGENT_SYSTEM_AVAILABLE:
                self.handle_nlp_parse_task_api()
            # Search API endpoints
            elif self.path == '/api/search' and SEARCH_AVAILABLE:
                self.handle_search_api()
            elif self.path == '/api/search/folders' and SEARCH_AVAILABLE:
                self.handle_search_create_folder_api()
            elif self.path == '/api/search/objects' and SEARCH_AVAILABLE:
                self.handle_search_add_object_api()
            # Prompts API endpoints
            elif self.path == '/api/prompts' and PROMPTS_AVAILABLE:
                self.handle_prompts_create_api()
            elif self.path == '/api/prompts/use' and PROMPTS_AVAILABLE:
                self.handle_prompts_use_api()
            elif self.path == '/api/prompts/search' and PROMPTS_AVAILABLE:
                self.handle_prompts_search_api()
            elif self.path.startswith('/api/prompts/') and PROMPTS_AVAILABLE:
                self.handle_prompts_update_api()
            # TODO API endpoints - specific routes first, then generic
            elif self.path == '/api/todos/users' and TODOS_AVAILABLE:
                self.handle_users_create_api()
            elif self.path == '/api/todos/lists' and TODOS_AVAILABLE:
                self.handle_lists_create_api()
            elif self.path == '/api/todos/tags' and TODOS_AVAILABLE:
                self.handle_tags_create_api()
            elif self.path == '/api/todos/parse' and TODOS_AVAILABLE:
                self.handle_todos_parse_api()
            elif self.path == '/api/todos/search' and TODOS_AVAILABLE:
                self.handle_todos_search_api()
            elif self.path == '/api/todos' and TODOS_AVAILABLE:
                self.handle_todos_create_api()
            elif self.path.startswith('/api/todos/lists/') and self.path.endswith('/share') and TODOS_AVAILABLE:
                self.handle_lists_share_api()
            elif self.path.startswith('/api/todos/') and self.path.endswith('/complete') and TODOS_AVAILABLE:
                self.handle_todos_complete_api()
            elif self.path.startswith('/api/todos/') and self.path.endswith('/archive') and TODOS_AVAILABLE:
                self.handle_todos_archive_api()
            elif self.path.startswith('/api/todos/') and TODOS_AVAILABLE:
                self.handle_todos_update_api()
            # Opportunities API endpoints
            elif self.path == '/api/opportunities':
                self.handle_opportunities_create_api()
            elif self.path.startswith('/api/opportunities/') and '/tasks' in self.path:
                # Handle POST /api/opportunities/{id}/tasks - create task
                self.handle_tasks_create_api()
            elif self.path.startswith('/api/tasks/'):
                # Handle POST /api/tasks/{id} - update task
                self.handle_task_update_api()
            elif self.path.startswith('/api/opportunities/'):
                # Could be PATCH/PUT for update
                self.handle_opportunities_update_api()
            else:
                self.send_error(404, f"API endpoint not found: {self.path}")
                
        except Exception as e:
            print(f"❌ Error handling POST {self.path}: {e}")
            self.send_error(500, f"Internal server error: {e}")
    
    def do_DELETE(self):
        """Handle DELETE requests for API endpoints."""
        try:
            # Parse request path for DELETE operations
            if self.path.startswith('/api/rag/documents/') and RAG_AVAILABLE:
                # Extract document ID from path
                document_id = self.path.split('/')[-1]
                self.handle_rag_document_delete_api(document_id)
            elif self.path.startswith('/api/cag/documents/') and CAG_AVAILABLE:
                # Extract document ID from path
                document_id = self.path.split('/')[-1]
                self.handle_cag_document_delete_api(document_id)
            elif self.path.startswith('/api/prompts/') and PROMPTS_AVAILABLE:
                # Extract prompt ID from path
                prompt_id = self.path.split('/')[-1]
                self.handle_prompts_delete_api(prompt_id)
            # Ollama agents API endpoints
            elif self.path.startswith('/api/ollama/agents/') and OLLAMA_AGENTS_AVAILABLE:
                self.handle_ollama_agent_detail_api()
            # TODO API endpoints - specific routes first, then generic
            elif self.path.startswith('/api/todos/users/') and TODOS_AVAILABLE:
                user_id = self.path.split('/')[-1]
                self.handle_users_delete_api(user_id)
            elif self.path.startswith('/api/todos/lists/') and '/share' in self.path and TODOS_AVAILABLE:
                list_id = self.path.split('/')[-2]
                self.handle_lists_unshare_api(list_id)
            elif self.path.startswith('/api/todos/lists/') and TODOS_AVAILABLE:
                list_id = self.path.split('/')[-1]
                self.handle_lists_delete_api(list_id)
            elif self.path.startswith('/api/todos/tags/') and TODOS_AVAILABLE:
                tag_id = self.path.split('/')[-1]
                self.handle_tags_delete_api(tag_id)
            elif self.path.startswith('/api/todos/') and TODOS_AVAILABLE:
                # Extract todo ID from path
                todo_id = self.path.split('/')[-1]
                self.handle_todos_delete_api(todo_id)
            elif self.path.startswith('/api/tasks/'):
                # Extract task ID from path
                task_id = self.path.split('/')[-1]
                self.handle_task_delete_api(task_id)
            elif self.path.startswith('/api/opportunities/'):
                # Extract opportunity ID from path
                opportunity_id = self.path.split('/')[-1]
                self.handle_opportunities_delete_api(opportunity_id)
            else:
                self.send_error(404, f"API endpoint not found: {self.path}")

        except Exception as e:
            print(f"❌ Error handling DELETE {self.path}: {e}")
            self.send_error(500, f"Internal server error: {e}")

    def do_PUT(self):
        """Handle PUT requests for API endpoints."""
        try:
            # Prompts API endpoints
            if self.path.startswith('/api/prompts/') and PROMPTS_AVAILABLE:
                self.handle_prompts_update_api()
            # Ollama agents API endpoints
            elif self.path.startswith('/api/ollama/agents/') and OLLAMA_AGENTS_AVAILABLE:
                self.handle_ollama_agent_detail_api()
            # TODO API endpoints - specific routes first, then generic
            elif self.path.startswith('/api/todos/users/') and TODOS_AVAILABLE:
                user_id = self.path.split('/')[-1]
                self.handle_users_update_api(user_id)
            elif self.path.startswith('/api/todos/lists/') and not '/share' in self.path and TODOS_AVAILABLE:
                list_id = self.path.split('/')[-1]
                self.handle_lists_update_api(list_id)
            elif self.path.startswith('/api/todos/tags/') and TODOS_AVAILABLE:
                tag_id = self.path.split('/')[-1]
                self.handle_tags_update_api(tag_id)
            elif self.path.startswith('/api/todos/') and TODOS_AVAILABLE:
                # Handle PUT /api/todos/{id} - Update individual todo
                self.handle_todos_update_api()
            else:
                self.send_error(404, f"API endpoint not found: {self.path}")
        except Exception as e:
            print(f"❌ Error handling PUT {self.path}: {e}")
            self.send_error(500, f"Internal server error: {e}")

    def handle_chat_api(self):
        """Handle chat API requests with MCP-style RAG tool integration."""
        handle_chat_route(self)

    def handle_mcp_tool_call(self, mcp_tool_call, model):
        """Handle MCP tool execution requests."""
        try:
            import asyncio
            import base64
            import tempfile
            from pathlib import Path

            # Add mcp-tools to path
            mcp_tools_path = os.path.join(os.path.dirname(__file__), 'mcp-tools')
            if mcp_tools_path not in sys.path:
                sys.path.insert(0, mcp_tools_path)

            tool_id = mcp_tool_call.get('tool_id', '')
            tool_name = mcp_tool_call.get('tool_name', '')
            inputs = mcp_tool_call.get('inputs', {})

            # Process file uploads: convert base64 content to temp files
            temp_files = []  # Track temporary files for cleanup
            processed_inputs = {}

            for key, value in inputs.items():
                if isinstance(value, dict) and 'filename' in value and 'content' in value:
                    # This is a file upload - save to temp file
                    try:
                        file_content = base64.b64decode(value['content'])
                        suffix = Path(value['filename']).suffix

                        # Create temp file with original extension
                        temp_file = tempfile.NamedTemporaryFile(
                            mode='wb',
                            suffix=suffix,
                            delete=False
                        )
                        temp_file.write(file_content)
                        temp_file.close()

                        temp_files.append(temp_file.name)
                        processed_inputs[key] = temp_file.name

                        debug_log(f"Saved uploaded file {value['filename']} to {temp_file.name}", "📁")
                    except Exception as e:
                        error_log(f"Failed to process uploaded file: {e}")
                        # Cleanup any temp files created so far
                        for temp_path in temp_files:
                            try:
                                os.unlink(temp_path)
                            except:
                                pass
                        self.send_json_response({
                            'error': f'Failed to process uploaded file: {str(e)}'
                        }, 400)
                        return
                else:
                    # Regular string value
                    processed_inputs[key] = value

            # Use processed inputs (with temp file paths) instead of raw inputs
            inputs = processed_inputs

            debug_log(f"Executing MCP tool: {tool_name} ({tool_id})", "🔧")
            debug_log(f"   Inputs: {inputs}")

            # Handle whitepaper-review tool
            if tool_id == 'whitepaper-review':
                # Get input parameters
                rubric_path = inputs.get('rubric_path', '')
                content_path = inputs.get('content_path', '')
                review_model = inputs.get('model', model)
                timeout_seconds = inputs.get('timeout_seconds', '30')

                if not rubric_path or not content_path:
                    self.send_json_response({
                        'error': 'Missing required parameters: rubric_path and content_path are required'
                    }, 400)
                    return

                # Verify files exist
                if not Path(rubric_path).exists():
                    self.send_json_response({
                        'error': f'Rubric file not found: {rubric_path}'
                    }, 400)
                    return

                if not Path(content_path).exists():
                    self.send_json_response({
                        'error': f'Content file not found: {content_path}'
                    }, 400)
                    return

                try:
                    # Import and execute the whitepaper review tool (refactored version)
                    from whitepaper_review_server import WhitePaperReviewServer

                    # Create async function to execute the tool
                    async def execute_review():
                        server = WhitePaperReviewServer(
                            ollama_url=os.getenv('OLLAMA_URL', 'http://localhost:11434'),
                            default_model=review_model,
                            default_timeout=int(timeout_seconds)
                        )

                        # Load documents using shared DocumentLoader
                        rubric_docs = await server.document_loader.load(rubric_path)
                        rubric_text = server.document_loader.combine_documents(rubric_docs, strategy="concatenate")

                        content_docs = await server.document_loader.load(content_path)
                        content_text = server.document_loader.combine_documents(content_docs, strategy="sections")

                        if not rubric_text:
                            return json.dumps({
                                "status": "error",
                                "message": f"Failed to load rubric from: {rubric_path}"
                            })

                        if not content_text:
                            return json.dumps({
                                "status": "error",
                                "message": f"Failed to load content from: {content_path}"
                            })

                        # Call the review method directly
                        result = await server._perform_review(
                            rubric=rubric_text,
                            content=content_text,
                            model=review_model,
                            timeout=int(timeout_seconds)
                        )

                        # Format the output
                        formatted_result = server._format_output(
                            result,
                            'markdown',
                            rubric_path,
                            content_path
                        )

                        return formatted_result

                    # Run the async function
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    try:
                        response_text = loop.run_until_complete(execute_review())
                        debug_log(f"MCP tool completed successfully", "✅")

                        self.send_json_response({
                            'response': response_text,
                            'model': review_model,
                            'mcp_tool': tool_name,
                            'status': 'success'
                        })
                    finally:
                        loop.close()

                except ImportError as e:
                    print(f"❌ Failed to import MCP tool: {e}")
                    self.send_json_response({
                        'error': f'Failed to import MCP tool: {str(e)}. Make sure required packages are installed (uv add mcp python-docx PyPDF2).'
                    }, 500)

            elif tool_id == 'powerpoint-template-fill':
                # Get input parameters
                template_path = inputs.get('template_path', '')
                documents_folder = inputs.get('documents_folder', '')
                output_path = inputs.get('output_path', 'filled_presentation.pptx') or 'filled_presentation.pptx'
                placeholder_style = inputs.get('placeholder_style', 'angle') or 'angle'
                extraction_strategy = inputs.get('extraction_strategy', 'llm_smart') or 'llm_smart'
                ppt_model = inputs.get('model', model)
                timeout_seconds = inputs.get('timeout_seconds', '180') or '180'
                preserve_formatting = inputs.get('preserve_formatting', True)

                debug_log(f"PowerPoint Tool - Received inputs:", "📄")
                debug_log(f"   template_path: {template_path}")
                debug_log(f"   documents_folder: {documents_folder}")
                debug_log(f"   output_path: {output_path}")

                if not template_path or not documents_folder:
                    print(f"❌ Missing parameters - template: {bool(template_path)}, folder: {bool(documents_folder)}")
                    self.send_json_response({
                        'error': 'Missing required parameters: template_path and documents_folder are required'
                    }, 400)
                    return

                # Verify template file exists
                if not Path(template_path).exists():
                    self.send_json_response({
                        'error': f'Template file not found: {template_path}'
                    }, 400)
                    return

                # Verify documents folder exists
                if not Path(documents_folder).exists():
                    self.send_json_response({
                        'error': f'Documents folder not found: {documents_folder}'
                    }, 400)
                    return

                try:
                    # Import and execute the PowerPoint template tool
                    from powerpoint_template_server import PowerPointTemplateServer

                    # Create async function to execute the tool
                    async def execute_fill():
                        server = PowerPointTemplateServer(
                            ollama_url=os.getenv('OLLAMA_URL', 'http://localhost:11434'),
                            default_model=ppt_model,
                            default_timeout=int(timeout_seconds)
                        )

                        # Import necessary modules for the call
                        import json

                        # Validate template
                        debug_log(f"Validating template: {template_path}", "📄")
                        validation = await server._validate_template(template_path)
                        debug_log(f"Template valid: {validation['slide_count']} slides, {validation['placeholder_count']} placeholders", "✅")

                        # Load template
                        from pptx import Presentation
                        presentation = Presentation(template_path)

                        # Extract placeholders
                        debug_log(f"Extracting placeholders with style: {placeholder_style}", "🔍")
                        placeholders = server._extract_placeholders(presentation, placeholder_style)
                        debug_log(f"Found {len(placeholders)} unique placeholders", "📋")

                        if not placeholders:
                            return json.dumps({
                                "status": "error",
                                "message": "No placeholders found in template"
                            })

                        # Load documents
                        debug_log(f"Loading documents from: {documents_folder}", "📂")
                        docs = await server.document_loader.load(documents_folder, recursive=False)
                        combined_docs = server.document_loader.combine_documents(docs, strategy="sections")
                        debug_log(f"Loaded {len(docs)} document(s)", "✅")

                        if not combined_docs:
                            return json.dumps({
                                "status": "error",
                                "message": f"No documents loaded from: {documents_folder}"
                            })

                        # Extract data using LLM
                        debug_log(f"Extracting data using strategy: {extraction_strategy}", "🤖")
                        extracted_data = await server._extract_data(
                            documents=combined_docs,
                            placeholders=list(placeholders.keys()),
                            strategy=extraction_strategy,
                            model=ppt_model,
                            timeout=int(timeout_seconds)
                        )
                        debug_log(f"Extracted {len(extracted_data)} values", "✅")

                        # Fill placeholders
                        debug_log(f"Filling placeholders in template", "✏️")
                        filled_count = server._fill_placeholders(
                            presentation=presentation,
                            placeholders=placeholders,
                            extracted_data=extracted_data,
                            preserve_formatting=preserve_formatting
                        )
                        debug_log(f"Filled {filled_count} placeholder(s)", "✅")

                        # Save output
                        output_full_path = Path(server.config.output_dir) / output_path
                        output_full_path.parent.mkdir(parents=True, exist_ok=True)
                        presentation.save(str(output_full_path))
                        debug_log(f"Saved to: {output_full_path}", "💾")

                        # Return result
                        return json.dumps({
                            "status": "success",
                            "output_file": str(output_full_path),
                            "placeholders_found": len(placeholders),
                            "placeholders_filled": filled_count,
                            "documents_processed": len(docs),
                            "extraction_strategy": extraction_strategy,
                            "model_used": ppt_model
                        }, indent=2)

                    # Run the async function
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    try:
                        response_text = loop.run_until_complete(execute_fill())
                        debug_log(f"MCP tool completed successfully", "✅")

                        # Parse the JSON response to display nicely
                        result = json.loads(response_text)

                        if result.get('status') == 'success':
                            # Format success message
                            message = f"""✅ **PowerPoint Template Filled Successfully**

**Output File:** `{result['output_file']}`
**Placeholders Found:** {result['placeholders_found']}
**Placeholders Filled:** {result['placeholders_filled']}
**Documents Processed:** {result['documents_processed']}
**Extraction Strategy:** {result['extraction_strategy']}
**Model Used:** {result['model_used']}

The filled PowerPoint presentation has been saved to `{result['output_file']}`.
"""
                            self.send_json_response({
                                'response': message,
                                'model': ppt_model,
                                'mcp_tool': tool_name,
                                'status': 'success',
                                'output_file': result['output_file']
                            })
                        else:
                            # Error in processing
                            self.send_json_response({
                                'response': result.get('message', 'Unknown error'),
                                'status': 'error'
                            })
                    finally:
                        loop.close()

                except ImportError as e:
                    print(f"❌ Failed to import MCP tool: {e}")
                    self.send_json_response({
                        'error': f'Failed to import MCP tool: {str(e)}. Make sure required packages are installed (uv add python-pptx).'
                    }, 500)

            else:
                # Unknown MCP tool
                self.send_json_response({
                    'error': f'Unknown MCP tool: {tool_id}'
                }, 400)

        except Exception as e:
            print(f"❌ MCP tool error: {e}")
            import traceback
            traceback.print_exc()
            self.send_json_response({
                'error': f'MCP tool error: {str(e)}'
            }, 500)
        finally:
            # Clean up temporary files
            for temp_path in temp_files:
                try:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
                        debug_log(f"Cleaned up temp file: {temp_path}", "🗑️")
                except Exception as cleanup_error:
                    error_log(f"Failed to clean up temp file {temp_path}: {cleanup_error}")

    def handle_models_api(self):
        """Handle models API requests to get available Ollama models."""
        handle_models_route(self)
    
    def handle_rag_status_api(self):
        """Handle RAG status API requests."""
        handle_rag_status_route(self)
    def handle_rag_search_api(self):
        """Handle RAG search API requests."""
        handle_rag_search_route(self)
    def handle_rag_query_api(self):
        """Handle RAG query API requests."""
        handle_rag_query_route(self)
    def handle_rag_ingest_api(self):
        """Handle RAG document ingestion API requests (supports both FormData files and JSON file paths)."""
        handle_rag_ingest_route(self)
    def _handle_file_upload(self):
        """Handle multipart form data file upload."""
        import cgi
        import tempfile

        try:
            # Parse multipart form data
            content_type = self.headers.get('Content-Type', '')
            
            # Create a temporary environment variable for CGI
            environ = {
                'REQUEST_METHOD': 'POST',
                'CONTENT_TYPE': content_type,
                'CONTENT_LENGTH': self.headers.get('Content-Length', '0')
            }
            
            form = cgi.FieldStorage(
                fp=self.rfile,
                headers=self.headers,
                environ=environ
            )
            
            # Get the uploaded file
            if 'file' not in form:
                self.send_json_response({'error': 'No file uploaded'}, 400)
                return
            
            file_item = form['file']
            if not file_item.filename:
                self.send_json_response({'error': 'No file selected'}, 400)
                return
            
            # Extract workspace/knowledge_base parameter
            knowledge_base = 'default'
            if 'knowledge_base' in form:
                knowledge_base = form['knowledge_base'].value
            
            debug_log(f"File upload for workspace: {knowledge_base}", "📤")
            
            # Create uploads directory if it doesn't exist
            uploads_dir = Path('uploads')
            uploads_dir.mkdir(exist_ok=True)
            
            # Save uploaded file
            file_path = uploads_dir / file_item.filename
            with open(file_path, 'wb') as f:
                f.write(file_item.file.read())
            
            debug_log(f"File uploaded: {file_item.filename} -> {file_path}", "📤")
            
            # Process the uploaded file with RAG system
            try:
                debug_log(f"Starting RAG ingestion for: {file_path} (workspace: {knowledge_base})", "🔄")
                ingest_result = handle_rag_ingest_request(str(file_path), knowledge_base)
                debug_log(f"RAG ingest: '{file_path}' -> {ingest_result.get('status', 'unknown')}", "📄")
                
                if ingest_result.get('status') == 'error':
                    print(f"❌ RAG ingestion failed: {ingest_result.get('message', 'Unknown error')}")
                    
            except Exception as ingest_error:
                print(f"❌ RAG ingestion exception for '{file_path}': {ingest_error}")
                import traceback
                traceback.print_exc()
                ingest_result = {
                    'status': 'error',
                    'message': f'RAG ingestion failed: {str(ingest_error)}'
                }
            
            # Clean up uploaded file after processing (optional)
            try:
                os.remove(file_path)
                debug_log(f"Cleaned up temporary file: {file_path}", "🧹")
            except OSError as cleanup_error:
                print(f"⚠️  Could not clean up file {file_path}: {cleanup_error}")
            
            self.send_json_response(ingest_result)
            
        except Exception as e:
            print(f"❌ File upload error: {e}")
            self.send_json_response({'error': f'File upload failed: {str(e)}'}, 500)
    
    def should_trigger_rag(self, message):
        """Determine if RAG should be used for this message."""
        # Keywords that suggest the user wants information from documents
        rag_triggers = [
            # Question words
            'what', 'how', 'when', 'where', 'why', 'which', 'who',
            # Information seeking
            'explain', 'describe', 'tell me about', 'information about',
            'details about', 'summary of', 'analyze', 'analysis',
            # Document-specific
            'document', 'file', 'data', 'spreadsheet', 'csv', 'table',
            'report', 'content', 'contains', 'mentions',
            # Search patterns
            'find', 'search', 'look for', 'show me', 'list',
            # Knowledge queries
            'know about', 'learn about', 'understand'
        ]
        
        message_lower = message.lower()
        
        # Check if message contains RAG trigger words
        for trigger in rag_triggers:
            if trigger in message_lower:
                return True
        
        # Check if it's a question (ends with ?)
        if message.strip().endswith('?'):
            return True
        
        return False
    
    def get_rag_enhanced_response(self, message, model, workspace='default'):
        """Get RAG-enhanced response using MCP-style tool integration."""
        try:
            # Use the RAG query endpoint
            rag_result = handle_rag_query_request(message, max_context=5, workspace=workspace)

            if rag_result['status'] == 'success':
                response_text = rag_result['answer']
                sources = rag_result.get('sources', [])
                model_used = rag_result.get('model_used', model)

                # Extract token usage from RAG result
                token_info = {
                    'prompt_tokens': rag_result.get('prompt_tokens', 0),
                    'completion_tokens': rag_result.get('completion_tokens', 0),
                    'total_tokens': rag_result.get('total_tokens', 0)
                }

                # Debug: Print sources information
                debug_log(f"RAG Debug - Sources count: {len(sources)}", "🔍")
                debug_log(f"RAG Debug - Sources type: {type(sources)}", "🔍")
                if sources:
                    debug_log(f"RAG Debug - First source: {sources[0]}", "🔍")
                else:
                    debug_log(f"RAG Debug - RAG result keys: {rag_result.keys()}", "🔍")

                # Add source attribution if sources exist
                if sources:
                    # Extract unique document names from sources
                    document_names = set()
                    for source in sources:
                        if 'metadata' in source and 'source' in source['metadata']:
                            file_path = source['metadata']['source']
                            # Extract just the filename from the full path
                            document_name = os.path.basename(file_path)
                            document_names.add(document_name)

                    if document_names:
                        doc_list = ', '.join(sorted(document_names))
                        source_info = f"\n\n📚 Sources consulted: {doc_list}"
                    else:
                        source_info = f"\n\n📚 Sources consulted: {len(sources)} document(s)"

                    response_text += source_info

                return response_text, model_used, sources, token_info
            else:
                # Fallback to standard Ollama if RAG fails
                print(f"⚠️ RAG failed, falling back to standard response: {rag_result.get('message', 'Unknown error')}")
                fallback_response, fallback_tokens = self.get_standard_ollama_response_with_tokens(message, model)
                return fallback_response, model, [], fallback_tokens

        except Exception as e:
            print(f"❌ RAG enhancement error: {e}")
            fallback_response, fallback_tokens = self.get_standard_ollama_response_with_tokens(message, model)
            return fallback_response, model, [], fallback_tokens
    
    def get_standard_ollama_response_with_tokens(self, message, model):
        """Get standard Ollama response with token usage information."""
        debug_log(f"Making fallback Ollama request with model: {model}", "📤")
        result = ollama_config.generate_response(model, message)

        if result['success']:
            response_text = result['data'].get('response', 'Sorry, I could not generate a response.')

            # Extract token usage
            prompt_tokens = result['data'].get('prompt_eval_count', 0)
            completion_tokens = result['data'].get('eval_count', 0)
            token_info = {
                'prompt_tokens': prompt_tokens,
                'completion_tokens': completion_tokens,
                'total_tokens': prompt_tokens + completion_tokens
            }

            return response_text, token_info
        else:
            return "Error: Unable to generate response", {'prompt_tokens': 0, 'completion_tokens': 0, 'total_tokens': 0}

    def get_standard_ollama_response(self, message, model):
        """Get standard Ollama response without RAG using robust configuration."""
        debug_log(f"Making fallback Ollama request with model: {model}", "📤")
        result = ollama_config.generate_response(model, message)
        
        if result['success']:
            response_text = result['data'].get('response', 'Sorry, I could not generate a response.')
            debug_log(f"Fallback Ollama response: {response_text[:50]}...", "🤖")
            return response_text
        else:
            error_msg = f"Ollama connection failed after {result['attempt']} attempts: {result['error']}"
            print(f"❌ Standard Ollama request failed: {error_msg}")
            return f"Sorry, I'm currently unable to process your request. {error_msg}"
    
    # Search API handlers
    def handle_search_api(self):
        """Handle universal search API requests."""
        handle_search_route(self)
    def handle_search_folders_api(self):
        """Handle folders listing API requests."""
        handle_search_folders_route(self)
    def handle_search_create_folder_api(self):
        """Handle folder creation API requests."""
        handle_search_create_folder_route(self)
    def handle_search_tags_api(self):
        """Handle tags listing API requests."""
        handle_search_tags_route(self)
    def handle_search_add_object_api(self):
        """Handle adding searchable objects API requests."""
        handle_search_add_object_route(self)
    def send_json_response(self, data, status_code=200):
        """Send a JSON response."""
        send_json_response_util(self, data, status_code)
    
    def do_HEAD(self):
        """Handle HEAD requests (like GET but without body)."""
        try:
            # Handle root path
            if self.path == '/':
                self.path = '/index.html'
            
            # Remove leading slash and resolve file path
            file_path = self.path.lstrip('/')
            full_path = os.path.join(os.getcwd(), file_path)
            
            # Check if file exists
            if not os.path.exists(full_path) or not os.path.isfile(full_path):
                self.send_error(404, f"File not found: {self.path}")
                return
            
            # Get file size and content type
            file_size = os.path.getsize(full_path)
            content_type = self.get_content_type(file_path)
            
            # Send headers only (no body for HEAD)
            self.send_response(200)
            self.send_header('Content-Type', content_type)
            self.send_header('Content-Length', str(file_size))
            self.send_header('Cache-Control', 'no-cache')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Access-Control-Allow-Methods', 'GET, HEAD, POST, OPTIONS')
            self.send_header('Access-Control-Allow-Headers', 'Content-Type')
            self.end_headers()
            
        except Exception as e:
            print(f"❌ Error handling HEAD {self.path}: {e}")
            self.send_error(500, f"Internal server error: {e}")
    
    def do_OPTIONS(self):
        """Handle OPTIONS requests for CORS preflight."""
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, HEAD, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()
    
    def get_content_type(self, file_path):
        """Determine content type based on file extension."""
        return get_content_type_util(file_path)

    def get_request_body(self):
        """Parse JSON request body safely."""
        return get_request_body_util(self)

    def get_query_params(self):
        """Parse URL query parameters."""
        from urllib.parse import urlparse, parse_qs
        parsed_url = urlparse(self.path)
        # parse_qs returns lists for each value, get first value for each param
        query_params = {k: v[0] for k, v in parse_qs(parsed_url.query).items()}
        return query_params

    def log_message(self, format, *args):
        """Override to provide cleaner logging."""
        return  # Disable default logging to reduce noise

    def handle_rag_documents_api(self):
        """Handle RAG documents listing API requests."""
        handle_rag_documents_route(self)
    def handle_rag_analytics_api(self):
        """Handle RAG analytics API requests."""
        handle_rag_analytics_route(self)
    def handle_rag_upload_api(self):
        """Handle RAG file upload API requests."""
        handle_rag_upload_route(self)
    def handle_rag_document_delete_api(self, document_id):
        """Handle RAG document deletion API requests."""
        handle_rag_document_delete_route(self, document_id)
    def handle_cag_status_api(self):
        """Handle CAG cache status API requests."""
        handle_cag_status_route(self)
    def handle_cag_load_api(self):
        """Handle CAG document loading API requests."""
        handle_cag_load_route(self)
    def handle_cag_clear_api(self):
        """Handle CAG cache clear API requests."""
        handle_cag_clear_route(self)
    def handle_cag_document_delete_api(self, document_id):
        """Handle CAG document deletion API requests."""
        handle_cag_document_delete_route(self, document_id)
    def handle_cag_query_api(self):
        """Handle CAG-enhanced query API requests."""
        handle_cag_query_route(self)
    def handle_knowledge_graph_api(self):
        """Handle knowledge graph visualization API requests using vector embeddings."""
        handle_knowledge_graph_route(self)
    def handle_performance_dashboard_api(self):
        """Handle performance dashboard API requests using real analytics data."""
        handle_performance_dashboard_route(self)
    def handle_search_results_api(self):
        """Handle search results explorer API requests using real RAG search."""
        handle_search_results_route(self)
    def handle_agents_api(self):
        """Handle /api/agents endpoint."""
        handle_agents_route(self)
    def handle_agents_metrics_api(self):
        """Handle /api/agents/metrics endpoint for real-time monitoring."""
        handle_agents_metrics_route(self)
    def handle_agents_detail_api(self):
        """Handle /api/agents/{agent_id} endpoint."""
        handle_agents_detail_route(self)

    # Ollama Agent handlers
    def handle_ollama_agents_api(self):
        """Handle /api/ollama/agents endpoint."""
        handle_ollama_agents_route(self)
    def handle_ollama_agent_detail_api(self):
        """Handle /api/ollama/agents/{agent_id} endpoint."""
        handle_ollama_agent_detail_route(self)
    def handle_ollama_agent_invoke_api(self):
        """Handle /api/ollama/agents/{agent_id}/invoke endpoint."""
        handle_ollama_agent_invoke_route(self)
    def handle_ollama_skills_api(self):
        """Handle /api/ollama/skills endpoint."""
        handle_ollama_skills_route(self)
    def handle_ollama_status_api(self):
        """Handle /api/ollama/status endpoint."""
        handle_ollama_status_route(self)

    def handle_mcp_servers_api(self):
        """Handle /api/mcp/servers endpoint."""
        handle_mcp_servers_route(self)
    def handle_mcp_server_action_api(self):
        """Handle /api/mcp/servers/{server_id}/action endpoint."""
        handle_mcp_server_action_route(self)
    def handle_nlp_parse_task_api(self):
        """Handle /api/nlp/parse-task endpoint."""
        handle_nlp_parse_task_route(self)
    def handle_nlp_capabilities_api(self):
        """Handle /api/nlp/capabilities endpoint."""
        handle_nlp_capabilities_route(self)
    def handle_mcp_metrics_api(self):
        """Handle /api/mcp/metrics endpoint."""
        handle_mcp_metrics_route(self)
    def handle_prompts_list_api(self):
        """Handle GET /api/prompts - List all prompts."""
        handle_prompts_list_route(self)
    def handle_prompts_create_api(self):
        """Handle POST /api/prompts - Create new prompt."""
        handle_prompts_create_route(self)
    def handle_prompts_detail_api(self):
        """Handle GET /api/prompts/{prompt_id} - Get prompt by ID."""
        handle_prompts_detail_route(self)
    def handle_prompts_update_api(self):
        """Handle PUT/POST /api/prompts/{prompt_id} - Update prompt."""
        handle_prompts_update_route(self)
    def handle_prompts_delete_api(self, prompt_id):
        """Handle DELETE /api/prompts/{prompt_id} - Delete prompt."""
        handle_prompts_delete_route(self, prompt_id)
    def handle_prompts_use_api(self):
        """Handle POST /api/prompts/use - Use a prompt (get content)."""
        handle_prompts_use_route(self)
    def handle_prompts_search_api(self):
        """Handle POST /api/prompts/search - Search prompts."""
        handle_prompts_search_route(self)

    # ========== Opportunities API Handlers ==========

    def handle_opportunities_list_api(self):
        """Handle GET /api/opportunities - List all opportunities."""
        handle_opportunities_list_route(self)

    def handle_opportunities_create_api(self):
        """Handle POST /api/opportunities - Create new opportunity."""
        handle_opportunities_create_route(self)

    def handle_opportunities_detail_api(self):
        """Handle GET /api/opportunities/{id} - Get opportunity details."""
        handle_opportunities_detail_route(self)

    def handle_opportunities_update_api(self):
        """Handle PUT/PATCH /api/opportunities/{id} - Update opportunity."""
        handle_opportunities_update_route(self)

    def handle_opportunities_delete_api(self, opportunity_id):
        """Handle DELETE /api/opportunities/{id} - Delete opportunity."""
        handle_opportunities_delete_route(self, opportunity_id)

    # ========== Tasks API Handlers ==========

    def handle_tasks_list_api(self):
        """Handle GET /api/opportunities/{id}/tasks - List tasks for opportunity."""
        # Extract opportunity ID from path: /api/opportunities/{id}/tasks
        path_parts = self.path.split('/')
        opportunity_id = path_parts[3]  # /api/opportunities/{id}/tasks
        handle_tasks_list_route(self, opportunity_id)

    def handle_tasks_create_api(self):
        """Handle POST /api/opportunities/{id}/tasks - Create new task."""
        # Extract opportunity ID from path: /api/opportunities/{id}/tasks
        path_parts = self.path.split('/')
        opportunity_id = path_parts[3]  # /api/opportunities/{id}/tasks
        handle_tasks_create_route(self, opportunity_id)

    def handle_task_get_api(self):
        """Handle GET /api/tasks/{id} - Get task details."""
        # Extract task ID from path: /api/tasks/{id}
        task_id = self.path.split('/')[-1]
        handle_task_get_route(self, task_id)

    def handle_task_update_api(self):
        """Handle POST /api/tasks/{id} - Update task."""
        # Extract task ID from path: /api/tasks/{id}
        task_id = self.path.split('/')[-1]
        handle_task_update_route(self, task_id)

    def handle_task_delete_api(self, task_id):
        """Handle DELETE /api/tasks/{id} - Delete task."""
        handle_task_delete_route(self, task_id)

    def handle_task_history_api(self):
        """Handle GET /api/tasks/{id}/history - Get task history."""
        # Extract task ID from path: /api/tasks/{id}/history
        task_id = self.path.split('/')[-2]  # /api/tasks/{id}/history
        handle_task_history_route(self, task_id)

    # TODO API handlers
    def handle_users_list_api(self):
        """Handle GET /api/todos/users - List users."""
        handle_users_list_route(self)

    def handle_users_create_api(self):
        """Handle POST /api/todos/users - Create user."""
        handle_users_create_route(self)

    def handle_users_detail_api(self):
        """Handle GET /api/todos/users/{id} - Get user details."""
        user_id = self.path.split('/')[-1]
        handle_users_detail_route(self, user_id)

    def handle_lists_list_api(self):
        """Handle GET /api/todos/lists - List todo lists."""
        handle_lists_list_route(self)

    def handle_lists_create_api(self):
        """Handle POST /api/todos/lists - Create todo list."""
        handle_lists_create_route(self)

    def handle_lists_detail_api(self):
        """Handle GET /api/todos/lists/{id} - Get list details."""
        list_id = self.path.split('/')[-1]
        handle_lists_detail_route(self, list_id)

    def handle_todos_list_api(self):
        """Handle GET /api/todos - List todos."""
        handle_todos_list_route(self)

    def handle_todos_create_api(self):
        """Handle POST /api/todos - Create todo."""
        handle_todos_create_route(self)

    def handle_todos_detail_api(self):
        """Handle GET /api/todos/{id} - Get todo details."""
        todo_id = self.path.split('/')[-1]
        handle_todos_detail_route(self, todo_id)

    def handle_todos_update_api(self):
        """Handle POST /api/todos/{id} - Update todo."""
        todo_id = self.path.split('/')[-1]
        handle_todos_update_route(self, todo_id)

    def handle_todos_delete_api(self, todo_id):
        """Handle DELETE /api/todos/{id} - Delete todo."""
        handle_todos_delete_route(self, todo_id)

    def handle_todos_complete_api(self):
        """Handle POST /api/todos/{id}/complete - Complete todo."""
        todo_id = self.path.split('/')[-2]  # /api/todos/{id}/complete
        handle_todos_complete_route(self, todo_id)

    def handle_todos_archive_api(self):
        """Handle POST /api/todos/{id}/archive - Archive todo."""
        todo_id = self.path.split('/')[-2]  # /api/todos/{id}/archive
        handle_todos_archive_route(self, todo_id)

    def handle_todos_today_api(self):
        """Handle GET /api/todos/today - Get today's todos."""
        handle_todos_today_route(self)

    def handle_todos_upcoming_api(self):
        """Handle GET /api/todos/upcoming - Get upcoming todos."""
        handle_todos_upcoming_route(self)

    def handle_todos_search_api(self):
        """Handle GET /api/todos/search - Search todos."""
        handle_todos_search_route(self)

    def handle_todos_parse_api(self):
        """Handle POST /api/todos/parse - Parse natural language input."""
        handle_todos_parse_route(self)

    def handle_tags_list_api(self):
        """Handle GET /api/todos/tags - List tags."""
        handle_tags_list_route(self)

    def handle_tags_create_api(self):
        """Handle POST /api/todos/tags - Create tag."""
        handle_tags_create_route(self)

    def handle_todos_history_api(self):
        """Handle GET /api/todos/{id}/history - Get todo history."""
        todo_id = self.path.split('/')[-2]  # /api/todos/{id}/history
        handle_todos_history_route(self, todo_id)

    # Phase 2 TODO API handlers

    def handle_users_update_api(self, user_id):
        """Handle PUT /api/todos/users/{id} - Update user."""
        handle_users_update_route(self, user_id)

    def handle_users_delete_api(self, user_id):
        """Handle DELETE /api/todos/users/{id} - Delete user."""
        handle_users_delete_route(self, user_id)

    def handle_lists_update_api(self, list_id):
        """Handle PUT /api/todos/lists/{id} - Update list."""
        handle_lists_update_route(self, list_id)

    def handle_lists_delete_api(self, list_id):
        """Handle DELETE /api/todos/lists/{id} - Delete list."""
        handle_lists_delete_route(self, list_id)

    def handle_lists_share_api(self):
        """Handle POST /api/todos/lists/{id}/share - Share list with user."""
        list_id = self.path.split('/')[-2]
        handle_lists_share_route(self, list_id)

    def handle_lists_unshare_api(self, list_id):
        """Handle DELETE /api/todos/lists/{id}/share - Unshare list."""
        handle_lists_unshare_route(self, list_id)

    def handle_lists_shares_api(self):
        """Handle GET /api/todos/lists/{id}/shares - Get list collaborators."""
        list_id = self.path.split('/')[-2]
        handle_lists_shares_route(self, list_id)

    def handle_shared_lists_api(self):
        """Handle GET /api/todos/shared - Get lists shared with user."""
        handle_shared_lists_route(self)

    def handle_tags_detail_api(self):
        """Handle GET /api/todos/tags/{id} - Get tag details."""
        tag_id = self.path.split('/')[-1]
        handle_tags_detail_route(self, tag_id)

    def handle_tags_update_api(self, tag_id):
        """Handle PUT /api/todos/tags/{id} - Update tag."""
        handle_tags_update_route(self, tag_id)

    def handle_tags_delete_api(self, tag_id):
        """Handle DELETE /api/todos/tags/{id} - Delete tag."""
        handle_tags_delete_route(self, tag_id)


def main():
    """Start the HTTP server on port 9090."""
    server_address = ('', 9090)
    httpd = HTTPServer(server_address, CustomHTTPRequestHandler)
    print(f'🚀 Server running on http://localhost:9090')
    print(f'📂 Serving files from: {os.getcwd()}')
    httpd.serve_forever()


if __name__ == '__main__':
    main()
