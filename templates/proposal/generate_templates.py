"""
Template Generator — Creates base proposal document templates.

Generates starter .docx, .xlsx, and .pptx files in the templates/proposal/
directory using python-docx, openpyxl, and python-pptx.

Run once to create/refresh the binary template files:
    uv run templates/proposal/generate_templates.py

Templates created:
    technical_volume.docx       — Technical Volume base
    management_volume.docx      — Management Volume base
    past_performance.docx       — Past Performance base
    cost_volume.xlsx            — Cost Volume Section B
    bid_no_bid.pptx             — Bid/No-Bid slide deck starter
    hotwash_report.docx         — Hotwash Report base
"""

import os
from pathlib import Path

# ---------------------------------------------------------------------------
# Directory setup
# ---------------------------------------------------------------------------

TEMPLATE_DIR = Path(__file__).parent


def _ensure_dir(path: Path) -> Path:
    """Create directory if it does not exist and return it."""
    path.mkdir(parents=True, exist_ok=True)
    return path


# ---------------------------------------------------------------------------
# Color palette (matches bid_no_bid_slide.py brand colors)
# ---------------------------------------------------------------------------

_COLOR_PRIMARY = os.getenv("BNB_COLOR_PRIMARY", "002060")    # Dark navy (no #)
_COLOR_ACCENT  = os.getenv("BNB_COLOR_ACCENT", "0070C0")    # Blue
_GRAY_LIGHT    = "F2F2F2"
_GRAY_MID      = "D9D9D9"
_WHITE         = "FFFFFF"


# ---------------------------------------------------------------------------
# DOCX templates
# ---------------------------------------------------------------------------

def _make_docx_template(output_path: Path, title: str, sections: list[str]) -> None:
    """
    Create a base Word document template with standard sections.

    Args:
        output_path: Where to write the .docx file.
        title: Document title for the cover heading.
        sections: List of section headings to include as placeholder pages.
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # --- Page margins ---
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.25)
        section.right_margin  = Inches(1.25)

    # --- Cover heading ---
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in heading.runs:
        run.font.color.rgb = RGBColor(0x00, 0x20, 0x60)

    doc.add_paragraph(
        "[ PROPOSAL TITLE ] | [ SOLICITATION NUMBER ] | [ DATE ]"
    ).alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # --- Table of Contents placeholder ---
    doc.add_heading("Table of Contents", level=1)
    doc.add_paragraph("[ Auto-generate Table of Contents here ]")
    doc.add_page_break()

    # --- Section placeholders ---
    for section_title in sections:
        doc.add_heading(section_title, level=1)
        doc.add_paragraph(
            f"[ {section_title.upper()} — Replace this placeholder with proposal content. ]"
        )
        doc.add_paragraph("")

    # --- Footer note ---
    doc.add_paragraph()
    note = doc.add_paragraph(
        "NOTICE: This document may contain proprietary information. "
        "Handle per contract requirements."
    )
    note.runs[0].font.size = Pt(8)
    note.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    doc.save(output_path)
    print(f"  Created: {output_path.relative_to(TEMPLATE_DIR.parent.parent)}")


def create_technical_volume() -> None:
    """Create Technical Volume base template."""
    _make_docx_template(
        TEMPLATE_DIR / "technical_volume.docx",
        "Technical Volume",
        [
            "1.0 Executive Summary",
            "2.0 Technical Approach",
            "2.1 Understanding of Requirements",
            "2.2 Technical Solution",
            "2.3 Management Approach",
            "3.0 Key Personnel",
            "4.0 Facilities and Equipment",
            "5.0 Risk Management",
        ],
    )


def create_management_volume() -> None:
    """Create Management Volume base template."""
    _make_docx_template(
        TEMPLATE_DIR / "management_volume.docx",
        "Management Volume",
        [
            "1.0 Management Approach",
            "2.0 Organizational Structure",
            "3.0 Key Personnel",
            "4.0 Transition Plan",
            "5.0 Quality Assurance Plan",
            "6.0 Subcontracting Plan",
        ],
    )


def create_past_performance() -> None:
    """Create Past Performance template."""
    _make_docx_template(
        TEMPLATE_DIR / "past_performance.docx",
        "Past Performance Volume",
        [
            "1.0 Introduction",
            "2.0 Relevant Past Performance Reference 1",
            "2.1 Contract Information",
            "2.2 Scope and Relevance",
            "2.3 Performance Summary",
            "2.4 Customer Reference",
            "3.0 Relevant Past Performance Reference 2",
            "4.0 Relevant Past Performance Reference 3",
        ],
    )


def create_hotwash_report() -> None:
    """Create Hotwash Report base template."""
    _make_docx_template(
        TEMPLATE_DIR / "hotwash_report.docx",
        "Proposal Hotwash Report",
        [
            "1.0 Proposal Overview",
            "2.0 Outcome Summary",
            "3.0 Government Debrief Notes",
            "4.0 Internal Assessment",
            "5.0 Lessons Learned",
            "6.0 Process Improvement Recommendations",
            "7.0 Action Items",
        ],
    )


# ---------------------------------------------------------------------------
# XLSX template — Cost Volume Section B
# ---------------------------------------------------------------------------

def create_cost_volume() -> None:
    """
    Create Cost Volume Section B Excel template.

    Includes:
        - Cover sheet with opportunity metadata fields
        - Labor sheet with CLIN/SLIN rows and wrap-rate columns
        - ODC sheet for Other Direct Costs
        - Summary sheet rolling up all CLINs
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ---- Helper styles ----
    header_fill  = PatternFill("solid", fgColor=_COLOR_PRIMARY)
    accent_fill  = PatternFill("solid", fgColor=_COLOR_ACCENT)
    gray_fill    = PatternFill("solid", fgColor=_GRAY_LIGHT)
    white_font   = Font(name="Calibri", bold=True, color=_WHITE, size=11)
    normal_font  = Font(name="Calibri", size=10)
    bold_font    = Font(name="Calibri", bold=True, size=10)
    thin_border  = Border(
        left=Side(style="thin"),  right=Side(style="thin"),
        top=Side(style="thin"),   bottom=Side(style="thin"),
    )

    def _hdr(ws, row: int, col: int, value: str, fill=None, font=None) -> None:
        cell = ws.cell(row=row, column=col, value=value)
        if fill:
            cell.fill = fill
        if font:
            cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thin_border

    def _set_col_widths(ws, widths: dict) -> None:
        for col_letter, width in widths.items():
            ws.column_dimensions[col_letter].width = width

    # ==============================
    # SHEET 1: Cover
    # ==============================
    cover = wb.active
    cover.title = "Cover"
    cover["A1"] = "COST VOLUME — SECTION B PRICE SCHEDULE"
    cover["A1"].font = Font(name="Calibri", bold=True, size=16, color=_COLOR_PRIMARY)
    cover.merge_cells("A1:D1")

    fields = [
        ("Solicitation Number", ""),
        ("Contract Title", ""),
        ("Offeror Name", ""),
        ("CAGE Code", ""),
        ("DUNS/UEI", ""),
        ("Proposal Date", ""),
        ("Period of Performance", ""),
        ("Contract Type", ""),
    ]
    for i, (label, _) in enumerate(fields, start=3):
        cover.cell(row=i, column=1, value=label).font = bold_font
        cover.cell(row=i, column=2, value="[ Fill in ]").font = normal_font
    _set_col_widths(cover, {"A": 28, "B": 40})

    # ==============================
    # SHEET 2: Labor
    # ==============================
    labor = wb.create_sheet("Labor")
    labor_headers = [
        "CLIN", "SLIN", "Labor Category", "Level",
        "Base Hours", "Option 1 Hours", "Option 2 Hours",
        "Fully-Loaded Rate ($/hr)",
        "Base Ext Cost", "Option 1 Ext Cost", "Option 2 Ext Cost",
        "Total Cost",
    ]
    for col_idx, hdr in enumerate(labor_headers, start=1):
        _hdr(labor, 1, col_idx, hdr, fill=header_fill, font=white_font)

    # Sample rows (3 CLIN placeholders)
    sample_rows = [
        ("0001", "0001AA", "Program Manager", "Senior", 1920, 1920, 1920, 225.00),
        ("0001", "0001AB", "Systems Engineer", "Mid",   1920, 1920, 1920, 165.00),
        ("0002", "0002AA", "Software Developer", "Senior", 1920, 1920, 1920, 185.00),
    ]
    for r_idx, row_data in enumerate(sample_rows, start=2):
        clin, slin, cat, lvl, bh, o1h, o2h, rate = row_data
        labor.cell(r_idx, 1, clin).font = normal_font
        labor.cell(r_idx, 2, slin).font = normal_font
        labor.cell(r_idx, 3, cat).font  = normal_font
        labor.cell(r_idx, 4, lvl).font  = normal_font
        labor.cell(r_idx, 5, bh).font   = normal_font
        labor.cell(r_idx, 6, o1h).font  = normal_font
        labor.cell(r_idx, 7, o2h).font  = normal_font
        labor.cell(r_idx, 8, rate).font = normal_font
        # Formulas for extended costs
        r = r_idx
        labor.cell(r, 9,  f"=E{r}*H{r}").font  = normal_font
        labor.cell(r, 10, f"=F{r}*H{r}").font  = normal_font
        labor.cell(r, 11, f"=G{r}*H{r}").font  = normal_font
        labor.cell(r, 12, f"=I{r}+J{r}+K{r}").font = bold_font

    _set_col_widths(labor, {
        "A": 8, "B": 10, "C": 28, "D": 10,
        "E": 12, "F": 14, "G": 14, "H": 22,
        "I": 16, "J": 18, "K": 18, "L": 14,
    })

    # ==============================
    # SHEET 3: ODC
    # ==============================
    odc = wb.create_sheet("ODC")
    odc_headers = ["CLIN", "Description", "Quantity", "Unit", "Unit Cost", "Total Cost", "Notes"]
    for col_idx, hdr in enumerate(odc_headers, start=1):
        _hdr(odc, 1, col_idx, hdr, fill=accent_fill, font=white_font)
    _set_col_widths(odc, {"A": 8, "B": 40, "C": 10, "D": 10, "E": 14, "F": 14, "G": 30})

    # ==============================
    # SHEET 4: Summary
    # ==============================
    summary = wb.create_sheet("Summary")
    summary["A1"] = "COST SUMMARY"
    summary["A1"].font = Font(name="Calibri", bold=True, size=14, color=_COLOR_PRIMARY)
    summary.merge_cells("A1:D1")

    sum_rows = [
        ("Direct Labor", "=SUM(Labor!L:L)", ""),
        ("Other Direct Costs (ODC)", "=SUM(ODC!F:F)", ""),
        ("Subtotal Direct Cost", "=B3+B4", ""),
        ("Overhead (Rate: ___%)", "=B5*0.0", "Enter overhead rate in formula"),
        ("G&A (Rate: ___%)",      "=B5*0.0", "Enter G&A rate in formula"),
        ("Total Cost",            "=B5+B6+B7", ""),
        ("Fee (Rate: ___%)",      "=B8*0.0", "Enter fee rate in formula"),
        ("Total Price",           "=B8+B9", ""),
    ]
    for r_idx, (label, formula, note) in enumerate(sum_rows, start=3):
        summary.cell(r_idx, 1, label).font = bold_font
        cell = summary.cell(r_idx, 2, formula)
        cell.font = normal_font
        cell.number_format = '"$"#,##0.00'
        summary.cell(r_idx, 3, note).font = Font(name="Calibri", size=9, italic=True)
    _set_col_widths(summary, {"A": 35, "B": 18, "C": 40})

    output_path = TEMPLATE_DIR / "cost_volume.xlsx"
    wb.save(output_path)
    print(f"  Created: {output_path.relative_to(TEMPLATE_DIR.parent.parent)}")


# ---------------------------------------------------------------------------
# PPTX template — Bid/No-Bid starter
# ---------------------------------------------------------------------------

def create_bid_no_bid_template() -> None:
    """
    Create a blank Bid/No-Bid PowerPoint template (6 slides).

    Slide layout matches bid_no_bid_slide.py for visual consistency.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Completely blank

    _primary = RGBColor(0x00, 0x20, 0x60)
    _accent  = RGBColor(0x00, 0x70, 0xC0)
    _white   = RGBColor(0xFF, 0xFF, 0xFF)
    _gray    = RGBColor(0xF2, 0xF2, 0xF2)

    def _add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(blank_layout)
        # Background box
        from pptx.util import Emu
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _primary
        bg.line.fill.background()

        tf = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
        p = tf.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = _white

        tf2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(1))
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0]
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xFF)

    def _add_content_slide(title: str, placeholder: str) -> None:
        slide = prs.slides.add_slide(blank_layout)
        # Header bar
        hdr = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _primary
        hdr.line.fill.background()

        tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.8))
        p = tf.text_frame.paragraphs[0]
        p.text = title
        run = p.runs[0]
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = _white

        # Body placeholder
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5))
        body.text_frame.word_wrap = True
        bp = body.text_frame.paragraphs[0]
        bp.text = placeholder
        brun = bp.runs[0]
        brun.font.size = Pt(14)
        brun.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    slide_defs = [
        ("BID / NO-BID ASSESSMENT", "[ Solicitation Number ] | [ Agency ] | [ Date ]"),
        ("Opportunity Summary", None),
        ("Shipley Scoring Matrix", None),
        ("Pwin Analysis", None),
        ("Win Themes & Risks", None),
        ("DECISION", None),
    ]

    for i, (title, subtitle) in enumerate(slide_defs):
        if i == 0:
            _add_title_slide(title, subtitle or "")
        else:
            _add_content_slide(title, f"[ {title} — replace with proposal-specific content ]")

    output_path = TEMPLATE_DIR / "bid_no_bid.pptx"
    prs.save(output_path)
    print(f"  Created: {output_path.relative_to(TEMPLATE_DIR.parent.parent)}")


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_all() -> None:
    """Generate all proposal document templates."""
    print("Generating proposal document templates...")
    _ensure_dir(TEMPLATE_DIR)

    create_technical_volume()
    create_management_volume()
    create_past_performance()
    create_hotwash_report()
    create_cost_volume()
    create_bid_no_bid_template()

    print(f"\nAll templates written to: {TEMPLATE_DIR}")


if __name__ == "__main__":
    generate_all()
