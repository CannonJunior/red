"""
Bid/No-Bid Decision Slide Generator.

Generates a 6-slide PPTX deck for leadership review using python-pptx.
Slide structure matches the bid-no-bid SKILL.md spec:

    Slide 1: Opportunity Summary
    Slide 2: Scoring Matrix (8-factor table with color coding)
    Slide 3: Pwin Analysis (competitive landscape)
    Slide 4: Win Themes & Discriminators
    Slide 5: Risks & Mitigations
    Slide 6: DECISION (BID / NO-BID / CONDITIONAL)

Branding colors are config-driven via environment variables:
    BNB_COLOR_PRIMARY   — hex RGB, default 002060 (dark blue)
    BNB_COLOR_BID       — hex for BID indicator, default 00B050 (green)
    BNB_COLOR_NOBID     — hex for NO-BID indicator, default FF0000 (red)
    BNB_COLOR_COND      — hex for CONDITIONAL, default FFC000 (amber)
"""

from __future__ import annotations

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from proposal.models import BidDecision, BidNoBidAssessment, Proposal

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def _rgb(hex_str: str) -> RGBColor:
    """
    Parse a 6-char hex string to RGBColor.

    Args:
        hex_str: 6-char hex without # prefix (e.g., "002060").

    Returns:
        RGBColor: pptx color object.
    """
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _colors() -> dict:
    """
    Load branding colors from env vars.

    Returns:
        dict: Keys: primary, bid, nobid, conditional, text_dark, text_light.
    """
    return {
        "primary":     _rgb(os.getenv("BNB_COLOR_PRIMARY", "002060")),
        "bid":         _rgb(os.getenv("BNB_COLOR_BID",     "00B050")),
        "nobid":       _rgb(os.getenv("BNB_COLOR_NOBID",   "FF0000")),
        "conditional": _rgb(os.getenv("BNB_COLOR_COND",    "FFC000")),
        "text_dark":   _rgb("1F1F1F"),
        "text_light":  _rgb("FFFFFF"),
        "row_alt":     _rgb("D9E1F2"),
        "score_high":  _rgb("C6EFCE"),   # green fill for scores 7-10
        "score_mid":   _rgb("FFEB9C"),   # amber fill for scores 4-6
        "score_low":   _rgb("FFC7CE"),   # red fill for scores 1-3
    }


def _decision_color(decision: BidDecision, colors: dict) -> RGBColor:
    """
    Return the brand color for a given decision.

    Args:
        decision: BidDecision enum value.
        colors: Colors dict from _colors().

    Returns:
        RGBColor: Color for the decision.
    """
    if decision == BidDecision.BID:
        return colors["bid"]
    elif decision == BidDecision.NO_BID:
        return colors["nobid"]
    else:
        return colors["conditional"]


# ---------------------------------------------------------------------------
# Slide building helpers
# ---------------------------------------------------------------------------

def _add_title_bar(slide, title: str, colors: dict, subtitle: str = "") -> None:
    """
    Add a colored title bar to the top of a slide.

    Args:
        slide: pptx Slide object.
        title: Title text.
        colors: Colors dict.
        subtitle: Optional subtitle text below the bar.
    """
    # Title bar background shape
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(10), Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["primary"]
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = colors["text_light"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.15), Inches(0.85),
            Inches(9.7), Inches(0.4),
        )
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.size = Pt(11)
        sub_run.font.italic = True
        sub_run.font.color.rgb = colors["text_light"]


def _add_text_body(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 13) -> None:
    """
    Add a plain text box to a slide.

    Args:
        slide: pptx Slide object.
        text: Content text (newlines are respected).
        left: Left position in inches.
        top: Top position in inches.
        width: Width in inches.
        height: Height in inches.
        font_size: Font size in points.
    """
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)


def _score_fill_color(score: Optional[float], colors: dict) -> RGBColor:
    """
    Map a 1-10 score to a traffic-light cell fill color.

    Args:
        score: Criterion score (1-10) or None.
        colors: Colors dict.

    Returns:
        RGBColor: Green (7-10), amber (4-6), or red (1-3). White if unscored.
    """
    if score is None:
        return _rgb("FFFFFF")
    if score >= 7:
        return colors["score_high"]
    elif score >= 4:
        return colors["score_mid"]
    return colors["score_low"]


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _slide1_opportunity_summary(prs: Presentation, proposal: Proposal, colors: dict) -> None:
    """Slide 1: Opportunity Summary — key solicitation facts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_title_bar(slide, "Opportunity Summary", colors)

    fields = [
        ("Solicitation",      proposal.solicitation_number or "TBD"),
        ("Title",             proposal.title),
        ("Agency",            proposal.agency or "TBD"),
        ("NAICS",             proposal.naics_code or "TBD"),
        ("Set-Aside",         (proposal.set_aside_type or "unknown").upper().replace("_", " ")),
        ("Estimated Value",   f"${proposal.estimated_value:,.0f}" if proposal.estimated_value else "TBD"),
        ("Proposal Due",      proposal.proposal_due_date or "TBD"),
        ("Recompete",         "Yes" if proposal.is_recompete else "No"),
        ("Incumbent",         proposal.incumbent or "Unknown"),
        ("Capture Manager",   proposal.capture_manager or "Unassigned"),
    ]

    col_left = [0.3, 5.1]
    row_top = 1.25
    row_h = 0.45
    for i, (label, value) in enumerate(fields):
        col = i % 2
        row = i // 2
        top = row_top + row * row_h
        left = col_left[col]

        lbl_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.2), Inches(0.38))
        lbl_tf = lbl_box.text_frame
        lbl_run = lbl_tf.paragraphs[0].add_run()
        lbl_run.text = f"{label}:"
        lbl_run.font.bold = True
        lbl_run.font.size = Pt(12)
        lbl_run.font.color.rgb = colors["primary"]

        val_box = slide.shapes.add_textbox(Inches(left + 2.2), Inches(top), Inches(2.3), Inches(0.38))
        val_tf = val_box.text_frame
        val_run = val_tf.paragraphs[0].add_run()
        val_run.text = str(value)
        val_run.font.size = Pt(12)
        val_run.font.color.rgb = colors["text_dark"]

    # Pwin badge if available
    if proposal.pwin_score is not None:
        pwin_pct = int(proposal.pwin_score * 100)
        badge = slide.shapes.add_shape(1, Inches(7.8), Inches(5.8), Inches(2.0), Inches(1.0))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors["primary"]
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run1 = badge_tf.paragraphs[0].add_run()
        run1.text = f"Pwin: {pwin_pct}%"
        run1.font.size = Pt(16)
        run1.font.bold = True
        run1.font.color.rgb = colors["text_light"]


def _slide2_scoring_matrix(prs: Presentation, assessment: BidNoBidAssessment, colors: dict) -> None:
    """Slide 2: Scoring Matrix — 8-factor table with traffic-light colors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    score = assessment.weighted_score()
    _add_title_bar(
        slide,
        "Bid/No-Bid Scoring Matrix",
        colors,
        subtitle=f"Weighted Score: {score:.1f}/100"
    )

    # Table: Name | Weight | Score | Weighted | Notes
    rows = len(assessment.criteria) + 1  # +1 for header
    cols = 5
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.3), Inches(1.2),
        Inches(9.4), Inches(5.2),
    ).table

    # Column widths
    col_widths = [2.8, 0.9, 0.7, 1.0, 4.0]
    for i, w in enumerate(col_widths):
        table.columns[i].width = int(Inches(w))

    headers = ["Criterion", "Weight", "Score", "Weighted", "Notes"]
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["primary"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = colors["text_light"]

    total_weight = sum(c.weight for c in assessment.criteria if c.score is not None)

    for ri, criterion in enumerate(assessment.criteria, start=1):
        weighted_val = ""
        if criterion.score is not None and total_weight > 0:
            contrib = (criterion.score * criterion.weight / total_weight) * 10
            weighted_val = f"{contrib:.1f}"

        row_bg = colors["row_alt"] if ri % 2 == 0 else _rgb("FFFFFF")
        score_bg = _score_fill_color(criterion.score, colors)

        values = [
            criterion.name,
            f"{criterion.weight:.2f}x",
            f"{criterion.score:.0f}" if criterion.score else "—",
            weighted_val or "—",
            criterion.notes or "",
        ]
        for ci, val in enumerate(values):
            cell = table.cell(ri, ci)
            if ci == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = score_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci in (1, 2, 3) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(10)
            run.font.color.rgb = colors["text_dark"]


def _slide3_pwin_analysis(prs: Presentation, proposal: Proposal, assessment: BidNoBidAssessment, colors: dict) -> None:
    """Slide 3: Pwin Analysis — competitive landscape and ghost analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Pwin Analysis & Competitive Landscape", colors)

    pwin_pct = int((proposal.pwin_score or 0) * 100)
    score = assessment.weighted_score()

    # Pwin gauge (text-based)
    gauge_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(1.0))
    gauge_tf = gauge_box.text_frame
    p = gauge_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Estimated Pwin: {pwin_pct}%"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = colors["primary"]

    # B/NB score comparison
    score_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.7), Inches(1.0))
    score_tf = score_box.text_frame
    sp = score_tf.paragraphs[0]
    sr = sp.add_run()
    sr.text = f"B/NB Score: {score:.1f}/100"
    sr.font.size = Pt(20)
    sr.font.bold = True
    sr.font.color.rgb = _decision_color(assessment.recommendation, colors)

    # Competitive factors from proposal metadata
    comp_text = "Competitive Factors:\n"
    if proposal.is_recompete:
        comp_text += f"  • Recompete — Incumbent: {proposal.incumbent or 'Unknown'}\n"
    if proposal.set_aside_type and proposal.set_aside_type.value != "unknown":
        sa = proposal.set_aside_type.value.upper().replace("_", " ")
        comp_text += f"  • Set-aside: {sa}\n"
    if proposal.naics_code:
        comp_text += f"  • NAICS: {proposal.naics_code} — {proposal.naics_description or ''}\n"
    if not comp_text.strip().endswith("Competitive Factors:"):
        _add_text_body(slide, comp_text, 0.3, 2.4, 9.4, 3.0, font_size=13)
    else:
        _add_text_body(slide, "No competitive intelligence captured yet.", 0.3, 2.4, 9.4, 1.0, font_size=13)


def _slide4_win_themes(prs: Presentation, assessment: BidNoBidAssessment, colors: dict) -> None:
    """Slide 4: Win Themes & Discriminators."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Win Themes & Discriminators", colors)

    # Win themes column
    wt_header = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.4))
    wth = wt_header.text_frame.paragraphs[0].add_run()
    wth.text = "Win Themes"
    wth.font.bold = True
    wth.font.size = Pt(14)
    wth.font.color.rgb = colors["primary"]

    themes_text = "\n".join(f"• {t}" for t in assessment.win_themes) if assessment.win_themes else "• (none recorded)"
    _add_text_body(slide, themes_text, 0.3, 1.65, 4.5, 5.0, font_size=12)

    # Discriminators column
    disc_header = slide.shapes.add_textbox(Inches(5.1), Inches(1.2), Inches(4.6), Inches(0.4))
    dh = disc_header.text_frame.paragraphs[0].add_run()
    dh.text = "Discriminators"
    dh.font.bold = True
    dh.font.size = Pt(14)
    dh.font.color.rgb = colors["primary"]

    disc_text = "\n".join(f"• {d}" for d in assessment.discriminators) if assessment.discriminators else "• (none recorded)"
    _add_text_body(slide, disc_text, 5.1, 1.65, 4.6, 5.0, font_size=12)

    # Divider line
    line = slide.shapes.add_shape(1, Inches(4.95), Inches(1.15), Inches(0.05), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = colors["primary"]
    line.line.fill.background()


def _slide5_risks(prs: Presentation, assessment: BidNoBidAssessment, colors: dict) -> None:
    """Slide 5: Risks & Mitigations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Risks & Mitigations", colors)

    risks = assessment.risks or ["(none recorded)"]
    mitigations = assessment.mitigations or []

    # Table: Risk | Mitigation
    rows = max(len(risks), 1) + 1
    table = slide.shapes.add_table(
        rows, 2,
        Inches(0.3), Inches(1.2),
        Inches(9.4), Inches(5.4),
    ).table

    table.columns[0].width = int(Inches(4.5))
    table.columns[1].width = int(Inches(4.9))

    for ci, hdr in enumerate(["Risk", "Mitigation"]):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["primary"]
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = colors["text_light"]

    for ri, risk_text in enumerate(risks, start=1):
        mit_text = mitigations[ri - 1] if ri - 1 < len(mitigations) else ""
        row_bg = colors["row_alt"] if ri % 2 == 0 else _rgb("FFFFFF")

        for ci, txt in enumerate([risk_text, mit_text]):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = txt
            run.font.size = Pt(11)
            run.font.color.rgb = colors["text_dark"]


def _slide6_decision(prs: Presentation, proposal: Proposal, assessment: BidNoBidAssessment, colors: dict) -> None:
    """Slide 6: DECISION — large-format decision announcement with rationale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "BID / NO-BID DECISION", colors)

    decision = assessment.final_decision
    dec_color = _decision_color(decision, colors)
    score = assessment.weighted_score()

    # Decision banner
    banner = slide.shapes.add_shape(
        1,  # rectangle
        Inches(1.5), Inches(1.5),
        Inches(7.0), Inches(2.0),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = dec_color
    banner.line.fill.background()
    btf = banner.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    brun = btf.paragraphs[0].add_run()
    brun.text = decision.value.upper().replace("_", " ")
    brun.font.size = Pt(48)
    brun.font.bold = True
    brun.font.color.rgb = colors["text_light"]

    # Score and date
    meta_box = slide.shapes.add_textbox(Inches(0.3), Inches(3.7), Inches(9.4), Inches(0.5))
    mp = meta_box.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = (
        f"Score: {score:.1f}/100  |  "
        f"Decision by: {assessment.decision_made_by or 'N/A'}  |  "
        f"Date: {(assessment.decision_date or '')[:10]}"
    )
    mr.font.size = Pt(12)
    mr.font.italic = True
    mr.font.color.rgb = colors["text_dark"]

    # Rationale
    if assessment.recommendation_rationale:
        _add_text_body(
            slide,
            f"Rationale:\n{assessment.recommendation_rationale}",
            0.3, 4.2, 9.4, 2.5,
            font_size=12,
        )

    # Solicitation reference footer
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.9), Inches(9.4), Inches(0.35))
    fp = footer_box.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"{proposal.solicitation_number}  |  {proposal.title[:80]}"
    fr.font.size = Pt(9)
    fr.font.color.rgb = _rgb("888888")


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def generate_bid_no_bid_slide(
    proposal: Proposal,
    assessment: BidNoBidAssessment,
    output_dir: Optional[str] = None,
) -> Path:
    """
    Generate a Bid/No-Bid PPTX deck for a proposal assessment.

    Saves the file to outputs/proposal/{solicitation}/ by default, following
    the project filename convention: {solicitation}_bid_no_bid_{YYYY-MM-DD}.pptx

    Args:
        proposal: Proposal model instance.
        assessment: Finalized BidNoBidAssessment.
        output_dir: Override output directory. Defaults to project standard.

    Returns:
        Path: Absolute path to the generated PPTX file.
    """
    colors = _colors()
    prs = Presentation()
    # Reason: Widescreen 16:9 is standard for modern executive presentations
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide1_opportunity_summary(prs, proposal, colors)
    _slide2_scoring_matrix(prs, assessment, colors)
    _slide3_pwin_analysis(prs, proposal, assessment, colors)
    _slide4_win_themes(prs, assessment, colors)
    _slide5_risks(prs, assessment, colors)
    _slide6_decision(prs, proposal, assessment, colors)

    # Build output path
    date_str = datetime.now().strftime("%Y-%m-%d")
    sol_id = (proposal.solicitation_number or proposal.id[:8]).replace("/", "-")

    if output_dir:
        out_path = Path(output_dir)
    else:
        # Reason: Follows project SKILL output convention from CLAUDE.md
        out_path = Path("outputs") / "proposal" / sol_id

    out_path.mkdir(parents=True, exist_ok=True)
    filename = out_path / f"{sol_id}_bid_no_bid_{date_str}.pptx"
    prs.save(str(filename))

    logger.info("Generated B/NB slide deck: %s", filename)
    return filename.resolve()
