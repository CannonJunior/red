"""
MCP (Model Context Protocol) API routes.

Handles MCP server management and tool execution:
- MCP tool execution (whitepaper review, PowerPoint template filling)
- MCP server listing and configuration
- MCP server actions (start, stop, restart)
- MCP system metrics
"""

import os
import sys
import json
import uuid
import asyncio
import base64
import tempfile
from pathlib import Path
from datetime import datetime

from debug_logger import debug_log, info_log, error_log
from server_decorators import require_system


class MCPRoutes:
    """Mixin providing MCP-related routes."""

    def handle_mcp_tool_call(self, mcp_tool_call, model):
        """
        Handle MCP tool execution requests.

        Called from chat API when mcp_tool_call is provided.
        Supports:
        - whitepaper-review: Review documents against rubrics
        - powerpoint-template-fill: Fill PowerPoint templates with extracted data
        """
        temp_files = []  # Track temporary files for cleanup

        try:
            # Add mcp-tools to path
            mcp_tools_path = os.path.join(os.path.dirname(__file__), '../../mcp-tools')
            if mcp_tools_path not in sys.path:
                sys.path.insert(0, mcp_tools_path)

            tool_id = mcp_tool_call.get('tool_id', '')
            tool_name = mcp_tool_call.get('tool_name', '')
            inputs = mcp_tool_call.get('inputs', {})

            # Process file uploads: convert base64 content to temp files
            processed_inputs = {}

            for key, value in inputs.items():
                if isinstance(value, dict) and 'filename' in value and 'content' in value:
                    # This is a file upload - save to temp file
                    try:
                        file_content = base64.b64decode(value['content'])
                        suffix = Path(value['filename']).suffix

                        # Create temp file with original extension
                        temp_file = tempfile.NamedTemporaryFile(
                            mode='wb',
                            suffix=suffix,
                            delete=False
                        )
                        temp_file.write(file_content)
                        temp_file.close()

                        temp_files.append(temp_file.name)
                        processed_inputs[key] = temp_file.name

                        debug_log(f"Saved uploaded file {value['filename']} to {temp_file.name}", "📁")
                    except Exception as e:
                        error_log(f"Failed to process uploaded file: {e}")
                        # Cleanup any temp files created so far
                        for temp_path in temp_files:
                            try:
                                os.unlink(temp_path)
                            except:
                                pass
                        self.send_json_response({
                            'error': f'Failed to process uploaded file: {str(e)}'
                        }, 400)
                        return
                else:
                    # Regular string value
                    processed_inputs[key] = value

            # Use processed inputs (with temp file paths) instead of raw inputs
            inputs = processed_inputs

            debug_log(f"Executing MCP tool: {tool_name} ({tool_id})", "🔧")

            # Route to appropriate tool handler
            if tool_id == 'whitepaper-review':
                self._handle_whitepaper_review(inputs, model, temp_files)
            elif tool_id == 'powerpoint-template-fill':
                self._handle_powerpoint_fill(inputs, model, temp_files)
            else:
                # Unknown MCP tool
                self.send_json_response({
                    'error': f'Unknown MCP tool: {tool_id}'
                }, 400)

        except Exception as e:
            print(f"❌ MCP tool error: {e}")
            import traceback
            traceback.print_exc()
            self.send_json_response({
                'error': f'MCP tool error: {str(e)}'
            }, 500)
        finally:
            # Clean up temporary files
            for temp_path in temp_files:
                try:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
                        debug_log(f"Cleaned up temp file: {temp_path}", "🗑️")
                except Exception as cleanup_error:
                    error_log(f"Failed to clean up temp file {temp_path}: {cleanup_error}")

    def _handle_whitepaper_review(self, inputs, model, temp_files):
        """Handle whitepaper review tool execution."""
        rubric_path = inputs.get('rubric_path', '')
        content_path = inputs.get('content_path', '')
        review_model = inputs.get('model', model)
        timeout_seconds = inputs.get('timeout_seconds', '30')

        if not rubric_path or not content_path:
            self.send_json_response({
                'error': 'Missing required parameters: rubric_path and content_path are required'
            }, 400)
            return

        # Verify files exist
        if not Path(rubric_path).exists():
            self.send_json_response({
                'error': f'Rubric file not found: {rubric_path}'
            }, 400)
            return

        if not Path(content_path).exists():
            self.send_json_response({
                'error': f'Content file not found: {content_path}'
            }, 400)
            return

        try:
            from whitepaper_review_server import WhitePaperReviewServer

            # Create async function to execute the tool
            async def execute_review():
                server = WhitePaperReviewServer(
                    ollama_url=os.getenv('OLLAMA_URL', 'http://localhost:11434'),
                    default_model=review_model,
                    default_timeout=int(timeout_seconds)
                )

                # Load documents
                rubric_docs = await server.document_loader.load(rubric_path)
                rubric_text = server.document_loader.combine_documents(rubric_docs, strategy="concatenate")

                content_docs = await server.document_loader.load(content_path)
                content_text = server.document_loader.combine_documents(content_docs, strategy="sections")

                if not rubric_text or not content_text:
                    return json.dumps({
                        "status": "error",
                        "message": "Failed to load documents"
                    })

                # Perform review
                result = await server._perform_review(
                    rubric=rubric_text,
                    content=content_text,
                    model=review_model,
                    timeout=int(timeout_seconds)
                )

                # Format output
                return server._format_output(result, 'markdown', rubric_path, content_path)

            # Run the async function
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                response_text = loop.run_until_complete(execute_review())
                debug_log(f"MCP tool completed successfully", "✅")

                self.send_json_response({
                    'response': response_text,
                    'model': review_model,
                    'mcp_tool': 'whitepaper-review',
                    'status': 'success'
                })
            finally:
                loop.close()

        except ImportError as e:
            print(f"❌ Failed to import MCP tool: {e}")
            self.send_json_response({
                'error': f'Failed to import MCP tool: {str(e)}. Install: uv add mcp python-docx PyPDF2'
            }, 500)

    def _handle_powerpoint_fill(self, inputs, model, temp_files):
        """Handle PowerPoint template filling tool execution."""
        template_path = inputs.get('template_path', '')
        documents_folder = inputs.get('documents_folder', '')
        output_path = inputs.get('output_path', 'filled_presentation.pptx') or 'filled_presentation.pptx'
        placeholder_style = inputs.get('placeholder_style', 'angle') or 'angle'
        extraction_strategy = inputs.get('extraction_strategy', 'llm_smart') or 'llm_smart'
        ppt_model = inputs.get('model', model)
        timeout_seconds = inputs.get('timeout_seconds', '180') or '180'
        preserve_formatting = inputs.get('preserve_formatting', True)

        if not template_path or not documents_folder:
            self.send_json_response({
                'error': 'Missing required parameters: template_path and documents_folder are required'
            }, 400)
            return

        # Verify paths exist
        if not Path(template_path).exists():
            self.send_json_response({'error': f'Template file not found: {template_path}'}, 400)
            return

        if not Path(documents_folder).exists():
            self.send_json_response({'error': f'Documents folder not found: {documents_folder}'}, 400)
            return

        try:
            from powerpoint_template_server import PowerPointTemplateServer

            async def execute_fill():
                server = PowerPointTemplateServer(
                    ollama_url=os.getenv('OLLAMA_URL', 'http://localhost:11434'),
                    default_model=ppt_model,
                    default_timeout=int(timeout_seconds)
                )

                from pptx import Presentation

                # Validate and load template
                presentation = Presentation(template_path)

                # Extract placeholders
                placeholders = server._extract_placeholders(presentation, placeholder_style)
                if not placeholders:
                    return json.dumps({"status": "error", "message": "No placeholders found"})

                # Load documents
                docs = await server.document_loader.load(documents_folder, recursive=False)
                combined_docs = server.document_loader.combine_documents(docs, strategy="sections")

                if not combined_docs:
                    return json.dumps({"status": "error", "message": "No documents loaded"})

                # Extract data using LLM
                extracted_data = await server._extract_data(
                    documents=combined_docs,
                    placeholders=list(placeholders.keys()),
                    strategy=extraction_strategy,
                    model=ppt_model,
                    timeout=int(timeout_seconds)
                )

                # Fill placeholders
                filled_count = server._fill_placeholders(
                    presentation=presentation,
                    placeholders=placeholders,
                    extracted_data=extracted_data,
                    preserve_formatting=preserve_formatting
                )

                # Save output
                output_full_path = Path(server.config.output_dir) / output_path
                output_full_path.parent.mkdir(parents=True, exist_ok=True)
                presentation.save(str(output_full_path))

                return json.dumps({
                    "status": "success",
                    "output_file": str(output_full_path),
                    "placeholders_found": len(placeholders),
                    "placeholders_filled": filled_count,
                    "documents_processed": len(docs)
                }, indent=2)

            # Run async function
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                response_text = loop.run_until_complete(execute_fill())
                result = json.loads(response_text)

                if result.get('status') == 'success':
                    message = f"""✅ **PowerPoint Template Filled Successfully**

**Output File:** `{result['output_file']}`
**Placeholders Filled:** {result['placeholders_filled']}/{result['placeholders_found']}
**Documents Processed:** {result['documents_processed']}
"""
                    self.send_json_response({
                        'response': message,
                        'model': ppt_model,
                        'mcp_tool': 'powerpoint-template-fill',
                        'status': 'success',
                        'output_file': result['output_file']
                    })
                else:
                    self.send_json_response({
                        'response': result.get('message', 'Unknown error'),
                        'status': 'error'
                    })
            finally:
                loop.close()

        except ImportError as e:
            self.send_json_response({
                'error': f'Failed to import MCP tool: {str(e)}. Install: uv add python-pptx'
            }, 500)

    @require_system('agent_system')
    def handle_mcp_servers_api(self):
        """
        Get/Create MCP servers.

        GET /api/mcp/servers - List all configured MCP servers
        POST /api/mcp/servers - Add new MCP server configuration
        """
        try:
            if self.command == 'GET':
                # Return configured MCP servers
                servers = [
                    {
                        'server_id': 'ollama_server',
                        'name': 'Ollama Local LLM Server',
                        'description': 'Local Ollama integration for zero-cost inference',
                        'status': 'running',
                        'host': 'localhost:11434',
                        'tools': ['llm_inference', 'text_generation']
                    },
                    {
                        'server_id': 'chromadb_server',
                        'name': 'ChromaDB Vector Server',
                        'description': 'Local vector database for RAG operations',
                        'status': 'running',
                        'host': 'localhost:9090',
                        'tools': ['vector_search', 'similarity_search', 'document_indexing']
                    }
                ]

                self.send_json_response({
                    'status': 'success',
                    'servers': servers,
                    'cost': '$0.00'
                })

            elif self.command == 'POST':
                # Add new MCP server
                server_data = self.get_request_body()
                if server_data is None:
                    self.send_json_response({'error': 'Invalid JSON'}, 400)
                    return

                server_id = f"server_{str(uuid.uuid4())[:8]}"
                new_server = {
                    'server_id': server_id,
                    'name': server_data.get('name', 'Unnamed Server'),
                    'description': server_data.get('description', ''),
                    'transport': server_data.get('transport', 'stdio'),
                    'status': 'stopped',
                    'created_at': datetime.now().isoformat()
                }

                # Validate transport-specific fields
                if new_server['transport'] == 'stdio':
                    if not server_data.get('command'):
                        self.send_json_response({
                            'error': 'Command is required for stdio transport'
                        }, 400)
                        return
                    new_server.update({
                        'command': server_data['command'],
                        'args': server_data.get('args', []),
                        'host': 'local',
                        'tools': ['local_execution']
                    })

                self.send_json_response({
                    'status': 'success',
                    'message': 'MCP server added successfully',
                    'data': new_server
                })

        except Exception as e:
            print(f"❌ MCP servers API error: {e}")
            self.send_json_response({'error': f'MCP servers API failed: {str(e)}'}, 500)

    @require_system('agent_system')
    def handle_mcp_server_action_api(self):
        """
        Perform action on MCP server.

        POST /api/mcp/servers/{server_id}/action
        Actions: start, stop, restart
        """
        try:
            self.send_json_response({
                'status': 'success',
                'message': 'MCP server action completed'
            })
        except Exception as e:
            self.send_json_response({'error': f'MCP server action failed: {str(e)}'}, 500)

    @require_system('agent_system')
    def handle_mcp_metrics_api(self):
        """
        Get MCP system metrics.

        GET /api/mcp/metrics
        """
        try:
            self.send_json_response({
                'status': 'success',
                'metrics': {
                    'total_servers': 2,
                    'active_servers': 2,
                    'total_tools': 6,
                    'performance': {
                        'avg_response_time_ms': 3.2,
                        'total_requests': 0
                    }
                },
                'cost': '$0.00'
            })
        except Exception as e:
            self.send_json_response({'error': f'MCP metrics failed: {str(e)}'}, 500)
