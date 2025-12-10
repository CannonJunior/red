#!/usr/bin/env python3
"""
PowerPoint Template Filling MCP Server

Fill PowerPoint templates with data extracted from document folders.
Uses shared infrastructure from common/ modules.
"""

import asyncio
import json
import os
import re
import sys
from pathlib import Path
from typing import Optional, Dict, Any, List, Set
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# MCP SDK imports
try:
    from mcp.server.fastmcp import FastMCP
    from mcp.types import Tool, TextContent
except ImportError:
    print("ERROR: MCP SDK not installed. Install with: uv add mcp")
    sys.exit(1)

# Shared infrastructure
try:
    from common.config import get_config
    from common.document_loader import DocumentLoader
    from common.ollama_client import OllamaClient
    from common.errors import DocumentLoadError, ExtractionError, TemplateError, MCPToolError
except ImportError:
    print("ERROR: Shared infrastructure not available")
    print("Make sure common/ modules are in mcp-tools/common/")
    sys.exit(1)


class PowerPointTemplateServer:
    """
    MCP server for filling PowerPoint templates with document data.

    Features:
    - Multi-format document extraction (.txt, .docx, .pdf)
    - LLM-powered intelligent data extraction
    - Flexible placeholder systems ({{mustache}}, [BRACKET])
    - Table support
    - Zero-cost local operation using Ollama
    - Uses shared infrastructure (DocumentLoader, OllamaClient, MCPToolConfig)
    """

    def __init__(
        self,
        ollama_url: Optional[str] = None,
        default_model: Optional[str] = None,
        default_timeout: Optional[int] = None,
        config=None
    ):
        """
        Initialize the PowerPoint Template Server.

        Args:
            ollama_url: Ollama API endpoint (uses config default if None)
            default_model: Default LLM model (uses config default if None)
            default_timeout: Default timeout in seconds (uses config default if None)
            config: MCPToolConfig instance (creates if None)
        """
        # Use shared configuration
        self.config = config or get_config()

        # Allow override of specific settings for backwards compatibility
        self.ollama_url = ollama_url or self.config.ollama_url
        self.default_model = default_model or self.config.default_model
        self.default_timeout = default_timeout or self.config.default_timeout

        # Initialize shared modules
        self.document_loader = DocumentLoader(config=self.config)
        self.ollama_client = OllamaClient(config=self.config)

        # Initialize MCP
        self.mcp = FastMCP("powerpoint-template-fill")

        # Register MCP tools
        self._register_tools()

    def _register_tools(self):
        """Register all MCP tools for PowerPoint template filling."""

        @self.mcp.tool()
        async def fill_powerpoint_template(
            template_path: str,
            documents_folder: str,
            output_path: Optional[str] = None,
            placeholder_style: str = "mustache",
            extraction_strategy: str = "llm_smart",
            model: Optional[str] = None,
            timeout_seconds: Optional[int] = None,
            preserve_formatting: bool = True
        ) -> str:
            """
            Fill PowerPoint template with data extracted from documents.

            Args:
                template_path: Path to .pptx template with placeholders
                documents_folder: Folder containing source documents
                output_path: Output file name (default: filled_presentation.pptx)
                placeholder_style: Placeholder format - 'mustache' or 'bracket'
                extraction_strategy: Data extraction method - 'llm_smart', 'keyword_match'
                model: LLM model to use (default: qwen2.5:3b)
                timeout_seconds: Maximum response time in seconds (default: 180)
                preserve_formatting: Attempt to preserve template text formatting

            Returns:
                JSON result with status and output file path
            """
            # Use defaults if not provided
            model = model or self.default_model
            timeout_seconds = timeout_seconds or 180
            output_path = output_path or "filled_presentation.pptx"

            try:
                # Validate template
                print(f"📄 Loading template: {template_path}")
                validation = await self._validate_template(template_path)
                print(f"✅ Template valid: {validation['slide_count']} slides, {validation['placeholder_count']} placeholders")

                # Load template
                presentation = Presentation(template_path)

                # Extract placeholders
                print(f"🔍 Extracting placeholders with style: {placeholder_style}")
                placeholders = self._extract_placeholders(presentation, placeholder_style)
                print(f"📋 Found {len(placeholders)} unique placeholders: {list(placeholders.keys())}")

                if not placeholders:
                    return json.dumps({
                        "status": "error",
                        "message": "No placeholders found in template",
                        "suggestions": [
                            f"Add {self._get_placeholder_example(placeholder_style)} markers to your template",
                            "Check placeholder_style parameter matches your template format"
                        ]
                    })

                # Load documents
                print(f"📂 Loading documents from: {documents_folder}")
                docs = await self.document_loader.load(documents_folder, recursive=False)
                combined_docs = self.document_loader.combine_documents(docs, strategy="sections")
                print(f"✅ Loaded {len(docs)} document(s), {len(combined_docs):,} chars")

                if not combined_docs:
                    return json.dumps({
                        "status": "error",
                        "message": f"No documents loaded from: {documents_folder}",
                        "suggestions": [
                            "Check the folder path is correct",
                            f"Ensure folder contains supported formats: {self.config.supported_document_formats}"
                        ]
                    })

                # Extract data using LLM
                print(f"🤖 Extracting data using strategy: {extraction_strategy}")
                extracted_data = await self._extract_data(
                    documents=combined_docs,
                    placeholders=list(placeholders.keys()),
                    strategy=extraction_strategy,
                    model=model,
                    timeout=timeout_seconds
                )
                print(f"✅ Extracted {len(extracted_data)} values")

                # Fill placeholders
                print(f"✏️  Filling placeholders in template")
                filled_count = self._fill_placeholders(
                    presentation=presentation,
                    placeholders=placeholders,
                    extracted_data=extracted_data,
                    preserve_formatting=preserve_formatting
                )
                print(f"✅ Filled {filled_count} placeholder(s)")

                # Save output
                output_full_path = Path(self.config.output_dir) / output_path
                output_full_path.parent.mkdir(parents=True, exist_ok=True)
                presentation.save(str(output_full_path))
                print(f"💾 Saved to: {output_full_path}")

                # Return result
                return json.dumps({
                    "status": "success",
                    "output_file": str(output_full_path),
                    "placeholders_found": len(placeholders),
                    "placeholders_filled": filled_count,
                    "documents_processed": len(docs),
                    "extraction_strategy": extraction_strategy,
                    "model_used": model,
                    "timestamp": datetime.now().isoformat()
                }, indent=2)

            except TemplateError as e:
                return json.dumps({
                    "status": "error",
                    "error_type": e.error_type,
                    "message": e.message,
                    "suggestions": e.suggestions,
                    "context": e.context
                })
            except DocumentLoadError as e:
                return json.dumps({
                    "status": "error",
                    "error_type": e.error_type,
                    "message": e.message,
                    "suggestions": e.suggestions
                })
            except ExtractionError as e:
                return json.dumps({
                    "status": "error",
                    "error_type": e.error_type,
                    "message": e.message,
                    "suggestions": e.suggestions
                })
            except Exception as e:
                return json.dumps({
                    "status": "error",
                    "message": f"Template filling failed: {str(e)}"
                })

    async def _validate_template(self, template_path: str) -> Dict[str, Any]:
        """
        Validate template before processing.

        Args:
            template_path: Path to template file

        Returns:
            Validation result dictionary

        Raises:
            TemplateError: If template is invalid
        """
        path = Path(template_path)

        # Check file exists
        if not path.exists():
            raise TemplateError(
                error_type="template_not_found",
                message=f"Template file not found: {template_path}",
                suggestions=["Check the file path is correct"]
            )

        # Check file extension
        if path.suffix.lower() != '.pptx':
            raise TemplateError(
                error_type="invalid_template_format",
                message=f"Template must be .pptx format, got: {path.suffix}",
                suggestions=["Ensure template is PowerPoint 2007+ format (.pptx)"]
            )

        # Check file size
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > self.config.max_file_size_mb:
            raise TemplateError(
                error_type="template_too_large",
                message=f"Template too large: {size_mb:.1f}MB (max {self.config.max_file_size_mb}MB)",
                suggestions=[f"Reduce template size below {self.config.max_file_size_mb}MB"]
            )

        # Try to load
        try:
            presentation = Presentation(template_path)
        except Exception as e:
            raise TemplateError(
                error_type="corrupt_template",
                message=f"Failed to load template: {str(e)}",
                suggestions=[
                    "Ensure file is a valid PowerPoint presentation",
                    "Try opening the file in PowerPoint and re-saving"
                ]
            )

        # Extract placeholders to count them
        placeholders = self._extract_placeholders(presentation, style="mustache")

        return {
            "valid": True,
            "slide_count": len(presentation.slides),
            "placeholder_count": len(placeholders),
            "placeholders": list(placeholders.keys()),
            "file_size_mb": round(size_mb, 2)
        }

    def _extract_placeholders(
        self,
        presentation: Presentation,
        style: str = "mustache"
    ) -> Dict[str, List[Dict[str, Any]]]:
        """
        Extract all placeholders from presentation.

        Args:
            presentation: python-pptx Presentation object
            style: Placeholder style ('mustache', 'bracket', or 'angle')

        Returns:
            Dictionary mapping placeholder names to list of locations
        """
        # Define regex patterns for different styles
        patterns = {
            "mustache": r'\{\{(\w+)\}\}',  # {{placeholder}}
            "bracket": r'\[(\w+)\]',  # [PLACEHOLDER]
            "angle": r'<([\w-]+)>'  # <placeholder> or <this-is-a-placeholder>
        }

        pattern = patterns.get(style, patterns["mustache"])
        placeholders = {}

        # Scan all slides
        for slide_idx, slide in enumerate(presentation.slides):
            # Scan all shapes
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape has text
                if not shape.has_text_frame:
                    continue

                # Extract text
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text

                # Find placeholders in text
                matches = re.findall(pattern, text)
                for match in matches:
                    placeholder_name = match.strip()
                    if placeholder_name not in placeholders:
                        placeholders[placeholder_name] = []

                    placeholders[placeholder_name].append({
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_type": str(shape.shape_type),
                        "original_text": text
                    })

        return placeholders

    def _get_placeholder_example(self, style: str) -> str:
        """Get example placeholder for given style."""
        examples = {
            "mustache": "{{company_name}}",
            "bracket": "[COMPANY_NAME]",
            "angle": "<company-name>"
        }
        return examples.get(style, "{{placeholder}}")

    async def _extract_data(
        self,
        documents: str,
        placeholders: List[str],
        strategy: str,
        model: str,
        timeout: int
    ) -> Dict[str, str]:
        """
        Extract data for placeholders from documents using specified strategy.

        Args:
            documents: Combined document text
            placeholders: List of placeholder names to extract
            strategy: Extraction strategy ('llm_smart' or 'keyword_match')
            model: LLM model name
            timeout: Timeout in seconds

        Returns:
            Dictionary mapping placeholder names to extracted values
        """
        if strategy == "llm_smart":
            return await self._extract_with_llm_smart(
                documents, placeholders, model, timeout
            )
        elif strategy == "keyword_match":
            return self._extract_with_keyword_match(documents, placeholders)
        else:
            raise ExtractionError(
                error_type="invalid_strategy",
                message=f"Unknown extraction strategy: {strategy}",
                suggestions=["Use 'llm_smart' or 'keyword_match'"]
            )

    async def _extract_with_llm_smart(
        self,
        documents: str,
        placeholders: List[str],
        model: str,
        timeout: int
    ) -> Dict[str, str]:
        """
        Extract data using LLM with smart prompt.

        Args:
            documents: Combined document text
            placeholders: List of placeholder names
            model: LLM model name
            timeout: Timeout in seconds

        Returns:
            Dictionary mapping placeholder names to extracted values
        """
        # Construct prompt
        placeholder_list = "\n".join([f"- {p}" for p in placeholders])

        prompt = f"""You are a data extraction assistant. Extract relevant content from the documents below for each placeholder.

PLACEHOLDERS TO FILL:
{placeholder_list}

DOCUMENTS:
{documents[:10000]}

For each placeholder, extract the most relevant content (1-3 sentences or a specific data point).
If a placeholder cannot be filled from the documents, respond with "NOT_FOUND".

Respond in JSON format:
{{
    "placeholder_name": "extracted value",
    ...
}}

Only include the JSON object, no other text."""

        try:
            # Call LLM
            response_text = await self.ollama_client.generate(
                prompt=prompt,
                model=model,
                timeout=timeout,
                temperature=0.3  # Lower temperature for more focused extraction
            )

            # Parse JSON response
            # Try to extract JSON from response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                extracted_data = json.loads(json_match.group())
            else:
                # Fallback: return empty dict
                print(f"⚠️  Could not parse JSON from LLM response, falling back to keyword matching")
                return self._extract_with_keyword_match(documents, placeholders)

            # Filter out NOT_FOUND values
            return {
                k: v for k, v in extracted_data.items()
                if v and v != "NOT_FOUND"
            }

        except Exception as e:
            print(f"⚠️  LLM extraction failed: {e}, falling back to keyword matching")
            return self._extract_with_keyword_match(documents, placeholders)

    def _extract_with_keyword_match(
        self,
        documents: str,
        placeholders: List[str]
    ) -> Dict[str, str]:
        """
        Extract data using simple keyword matching.

        Args:
            documents: Combined document text
            placeholders: List of placeholder names

        Returns:
            Dictionary mapping placeholder names to extracted values
        """
        extracted = {}

        for placeholder in placeholders:
            # Convert placeholder to search keywords
            # e.g., "company_name" -> "company name"
            keywords = placeholder.replace("_", " ").lower()

            # Search for keywords in documents
            sentences = documents.split(". ")
            for sentence in sentences:
                if keywords in sentence.lower():
                    # Return the sentence containing the keyword
                    extracted[placeholder] = sentence.strip()
                    break

        return extracted

    def _fill_placeholders(
        self,
        presentation: Presentation,
        placeholders: Dict[str, List[Dict[str, Any]]],
        extracted_data: Dict[str, str],
        preserve_formatting: bool
    ) -> int:
        """
        Fill placeholders in presentation with extracted data.

        Args:
            presentation: python-pptx Presentation object
            placeholders: Dictionary of placeholder locations
            extracted_data: Dictionary of extracted values
            preserve_formatting: Whether to preserve text formatting

        Returns:
            Number of placeholders filled
        """
        filled_count = 0

        for placeholder_name, locations in placeholders.items():
            # Get extracted value
            value = extracted_data.get(placeholder_name)

            if not value:
                print(f"⚠️  No data found for placeholder: {placeholder_name}")
                continue

            # Fill all locations of this placeholder
            for location in locations:
                slide_idx = location["slide_index"]
                shape_idx = location["shape_index"]

                try:
                    slide = presentation.slides[slide_idx]
                    shape = slide.shapes[shape_idx]

                    if shape.has_text_frame:
                        # Replace placeholder in text
                        original_text = location["original_text"]

                        # Build replacement pattern based on placeholder style
                        patterns = [
                            (r'\{\{' + re.escape(placeholder_name) + r'\}\}', value),  # {{placeholder}}
                            (r'\[' + re.escape(placeholder_name) + r'\]', value),  # [PLACEHOLDER]
                            (r'<' + re.escape(placeholder_name) + r'>', value)  # <placeholder>
                        ]

                        new_text = original_text
                        for pattern, replacement in patterns:
                            new_text = re.sub(pattern, replacement, new_text)

                        if preserve_formatting:
                            # Try to preserve formatting by only replacing text in runs
                            self._replace_text_preserve_formatting(shape.text_frame, original_text, new_text)
                        else:
                            # Simple replacement (loses formatting)
                            shape.text = new_text

                        filled_count += 1

                except Exception as e:
                    print(f"⚠️  Failed to fill placeholder {placeholder_name} at slide {slide_idx}: {e}")

        return filled_count

    def _replace_text_preserve_formatting(self, text_frame, old_text, new_text):
        """
        Replace text while attempting to preserve formatting.

        Args:
            text_frame: python-pptx TextFrame object
            old_text: Original text
            new_text: New text to replace with
        """
        # This is a simplified version - full preservation is complex
        # For now, just replace the text in the first run
        if text_frame.paragraphs:
            for paragraph in text_frame.paragraphs:
                if paragraph.runs:
                    # Replace in first run
                    paragraph.runs[0].text = new_text
                    # Clear other runs
                    for run in paragraph.runs[1:]:
                        run.text = ""
                    return

        # Fallback: set text_frame text directly
        text_frame.text = new_text

    async def run(self):
        """Run the MCP server."""
        print("🚀 Starting PowerPoint Template Fill MCP Server...")
        print(f"📋 Ollama URL: {self.ollama_url}")
        print(f"🤖 Default Model: {self.default_model}")
        print(f"⏱️  Default Timeout: {self.default_timeout}s")
        print(f"💾 Cache Enabled: {self.config.enable_cache}")
        print(f"📂 Output Dir: {self.config.output_dir}")
        print("✅ MCP Server ready for connections")

        await self.mcp.run()


async def main():
    """Main entry point for the MCP server."""
    # Read configuration from environment variables (for backwards compatibility)
    ollama_url = os.getenv('OLLAMA_URL')
    default_model = os.getenv('DEFAULT_MODEL')
    default_timeout = int(os.getenv('DEFAULT_TIMEOUT', '0')) or None

    # Create and run server (will use MCPToolConfig defaults if env vars not set)
    server = PowerPointTemplateServer(
        ollama_url=ollama_url,
        default_model=default_model,
        default_timeout=default_timeout
    )

    await server.run()


if __name__ == '__main__':
    # Run the server
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("\n👋 PowerPoint Template Fill MCP Server stopped")
    except Exception as e:
        print(f"❌ Server error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
